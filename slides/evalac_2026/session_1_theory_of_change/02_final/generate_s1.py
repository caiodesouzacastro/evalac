"""Generate Session 1 PPTX — EVALAC I: Introduction to M&E and Theory of Change."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x75, 0x75, 0x75)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
EXERCISE_BG = RGBColor(0xE3, 0xF2, 0xFD)
EXERCISE_ACCENT = RGBColor(0x15, 0x65, 0xC0)
DEBRIEF_BG = RGBColor(0xFE, 0xF3, 0xE0)
TRANSITION_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
REF_BG = RGBColor(0xEC, 0xEF, 0xF1)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color, x=0, y=0, w=None, h=None):
    w = w or prs.slide_width
    h = h or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(4), font_name="Calibri", italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_run(p, text, size=18, bold=False, color=DARK_GRAY, italic=False, font_name="Calibri"):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = font_name
    return run


def make_title_slide(title, subtitle="", slide_num=None, bg_color=DARK_BLUE, title_color=WHITE, sub_color=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, bg_color)

    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    tb = add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(2.5))
    set_text(tb.text_frame, title, size=36, bold=True, color=title_color, alignment=PP_ALIGN.LEFT)

    if subtitle:
        add_paragraph(tb.text_frame, subtitle, size=20, color=sub_color, space_before=Pt(16))

    if slide_num is not None:
        sn = add_text_box(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
        set_text(sn.text_frame, str(slide_num), size=12, color=RGBColor(0xAA, 0xAA, 0xAA), alignment=PP_ALIGN.RIGHT)

    return slide


def make_content_slide(title, bullet_items, slide_num=None, note_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    # top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.background()

    # title
    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9))
    set_text(tb.text_frame, title, size=28, bold=True, color=DARK_BLUE)

    # content
    cb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    cb.text_frame.word_wrap = True
    first = True
    for item in bullet_items:
        if first:
            p = cb.text_frame.paragraphs[0]
            first = False
        else:
            p = cb.text_frame.add_paragraph()

        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}

        p.text = text
        p.font.size = Pt(opts.get("size", 18))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK_GRAY)
        p.font.italic = opts.get("italic", False)
        p.font.name = "Calibri"
        p.space_before = Pt(opts.get("space_before", 6))
        p.space_after = Pt(opts.get("space_after", 4))
        if opts.get("indent"):
            p.level = opts["indent"]

    if slide_num is not None:
        sn = add_text_box(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
        set_text(sn.text_frame, str(slide_num), size=12, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

    if note_text:
        slide.notes_slide.notes_text_frame.text = note_text

    return slide


def make_exercise_slide(title, instructions, slide_num=None, time_label="", note_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EXERCISE_BG)

    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = EXERCISE_ACCENT
    bar.line.fill.background()

    # gear icon label
    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    add_run(p, "EXERCISE  ", size=14, bold=True, color=EXERCISE_ACCENT)
    if time_label:
        add_run(p, time_label, size=14, color=MED_GRAY)

    # title
    add_paragraph(tb.text_frame, title, size=26, bold=True, color=DARK_BLUE, space_before=Pt(4))

    cb = add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
    cb.text_frame.word_wrap = True
    first = True
    for item in instructions:
        if first:
            p = cb.text_frame.paragraphs[0]
            first = False
        else:
            p = cb.text_frame.add_paragraph()

        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}

        p.text = text
        p.font.size = Pt(opts.get("size", 18))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK_GRAY)
        p.font.italic = opts.get("italic", False)
        p.font.name = "Calibri"
        p.space_before = Pt(opts.get("space_before", 6))
        p.space_after = Pt(opts.get("space_after", 4))

    if slide_num:
        sn = add_text_box(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
        set_text(sn.text_frame, str(slide_num), size=12, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

    if note_text:
        slide.notes_slide.notes_text_frame.text = note_text

    return slide


def make_debrief_slide(title, points, slide_num=None, note_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEBRIEF_BG)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    add_run(p, "DEBRIEF  ", size=14, bold=True, color=ACCENT_ORANGE)

    add_paragraph(tb.text_frame, title, size=26, bold=True, color=DARK_BLUE, space_before=Pt(4))

    cb = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.3))
    cb.text_frame.word_wrap = True
    first = True
    for item in points:
        if first:
            p = cb.text_frame.paragraphs[0]
            first = False
        else:
            p = cb.text_frame.add_paragraph()
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p.text = text
        p.font.size = Pt(opts.get("size", 18))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK_GRAY)
        p.font.italic = opts.get("italic", False)
        p.font.name = "Calibri"
        p.space_before = Pt(opts.get("space_before", 6))
        p.space_after = Pt(opts.get("space_after", 4))

    if slide_num:
        sn = add_text_box(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
        set_text(sn.text_frame, str(slide_num), size=12, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

    if note_text:
        slide.notes_slide.notes_text_frame.text = note_text

    return slide


def make_section_slide(block_label, block_title, time_info="", slide_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MID_BLUE)

    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(3))
    p = tb.text_frame.paragraphs[0]
    add_run(p, block_label, size=16, bold=True, color=RGBColor(0xBB, 0xDE, 0xFB))
    add_paragraph(tb.text_frame, block_title, size=36, bold=True, color=WHITE, space_before=Pt(12))
    if time_info:
        add_paragraph(tb.text_frame, time_info, size=16, color=RGBColor(0xBB, 0xDE, 0xFB), space_before=Pt(12))

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.8), Inches(4), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()

    if slide_num:
        sn = add_text_box(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
        set_text(sn.text_frame, str(slide_num), size=12, color=RGBColor(0x88, 0x88, 0xBB), alignment=PP_ALIGN.RIGHT)

    return slide


def make_table_slide(title, headers, rows, slide_num=None, note_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MID_BLUE
    bar.line.fill.background()

    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8))
    set_text(tb.text_frame, title, size=28, bold=True, color=DARK_BLUE)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_w = Inches(11.5)
    tbl_h = Inches(min(5.2, 0.5 + 0.5 * n_rows))
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(0.8), Inches(1.4), tbl_w, tbl_h).table

    col_w = tbl_w // n_cols
    for i in range(n_cols):
        tbl.columns[i].width = col_w

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = MID_BLUE

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE

    if slide_num:
        sn = add_text_box(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
        set_text(sn.text_frame, str(slide_num), size=12, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

    if note_text:
        slide.notes_slide.notes_text_frame.text = note_text

    return slide


# ═══════════════════════════════════════════════════════════════════
# SLIDE GENERATION
# ═══════════════════════════════════════════════════════════════════

# ── SLIDE 1: Welcome ──
s = make_title_slide(
    "EVALAC I — Introduction to M&E\nand Theory of Change",
    "CLEAR-LAB  ·  IDB  ·  OVE",
    slide_num=1
)
s.notes_slide.notes_text_frame.text = (
    "Brief welcome (2 min). Logistics (1 min). "
    "Invite participants to introduce themselves with one sentence: name, country, current M&E role (2 min). "
    "Keep this short — the session is dense."
)

# ── SLIDE 2: How this week works ──
make_content_slide(
    "How This Week Works: 5 Sessions, One Story",
    [
        ("Session 1: How is change intended to happen?  (Theory of Change)", {"bold": True, "color": LIGHT_BLUE}),
        ("Session 2: How do we measure it?  (Indicators & Results Matrix)", {}),
        ("Session 3: What qualitative evidence do we need?  (Qualitative Methods)", {}),
        ("Session 4: What quantitative evidence do we need?  (Quantitative Methods)", {}),
        ("Session 5: How strong is our contribution claim?  (Contribution Analysis)", {}),
        ("", {}),
        ("All 5 sessions use the same case study. We build the ToC today and carry it forward all week.", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=2,
    note_text="Emphasize: all 5 sessions use the same case study. Session 2 picks up exactly where Session 1 ends."
)

# ── SLIDE 3: Learning objectives ──
make_content_slide(
    "What We Will Learn Today: Session 1 Objectives",
    [
        ("1.  Frame complex, multi-actor interventions as systems of change", {}),
        ("2.  Explain why theory-based evaluation is essential and position the ToC as the starting point", {}),
        ("3.  Describe how change is expected to happen: results, causal mechanisms, assumptions, and external influences", {}),
        ("4.  Distinguish monitoring from evaluation and explain the role of the ToC in both", {}),
        ("5.  Formulate evaluation questions that translate the ToC into evaluative inquiry", {}),
    ],
    slide_num=3,
    note_text="Read objectives aloud; do not explain at length. Return to this slide at the end as a check."
)

# ── SLIDE 4: Diagnostic survey ──
make_content_slide(
    "Diagnostic Survey: 3 Quick Questions",
    [
        ("Q1: Have you ever built or revised a Theory of Change?", {"bold": True}),
        ("    Yes  /  No  /  Tried, not sure it worked", {"size": 16, "color": MED_GRAY}),
        ("", {"size": 8}),
        ("Q2: What evaluation approach have you used most?", {"bold": True}),
        ("    RCT  /  Quasi-experimental  /  Process tracing  /  Contribution analysis  /  Other", {"size": 16, "color": MED_GRAY}),
        ("", {"size": 8}),
        ("Q3: What is your biggest M&E challenge right now?", {"bold": True}),
        ("    (1 sentence in chat or sticky note)", {"size": 16, "color": MED_GRAY}),
    ],
    slide_num=4,
    note_text="Display results live if using Miro/Mentimeter. Briefly name the distribution. Do not turn Q3 into a discussion — collect and set aside. Use Q3 responses as real-world anchors."
)

# ── SLIDE 5: Meet the case ──
s = make_content_slide(
    "Meet the Case: Creciendo Juntos",
    [
        ("THE CASE: Creciendo Juntos (Growing Together)", {"bold": True, "color": DARK_BLUE, "size": 22}),
        ("", {"size": 8}),
        ("A five-year, government-run early childhood development program serving 180 municipalities in a LAC country.", {}),
        ("Target group: children aged 0\u20133 from the poorest 40% of families.", {}),
        ("Reached group: ~85,000 caregivers (mothers, fathers, grandmothers).", {}),
        ("", {"size": 8}),
        ("Four components:", {"bold": True}),
        ("  \u2022  Parenting group sessions (36/year)     \u2022  Monthly home visits", {}),
        ("  \u2022  Growth monitoring at health centers   \u2022  Early learning centers", {}),
        ("", {"size": 8}),
        ("Budget: USD 285M over 5 years (~USD 1,140 per child/year). Delivered by 750 community facilitators.", {}),
        ("", {"size": 8}),
        ("Central question: Does participation by caregivers improve child development outcomes by age 3?", {"bold": True, "italic": True, "color": ACCENT_ORANGE}),
    ],
    slide_num=5,
    note_text="Read the vignette aloud (2 min). Pause on one question: 'Looking at this program — if it works, how exactly does that happen? What has to change?' Hold the question; do not answer yet."
)

# ── BLOCK 2 SECTION DIVIDER ──
make_section_slide("BLOCK 2", "Interventions and Their Results", "~55 minutes  \u00b7  Slides 6\u201313", slide_num=None)

# ── SLIDE 6: What is a public intervention? ──
make_table_slide(
    "What Is a Public Intervention?",
    ["INPUTS — What we invest", "ACTIVITIES — What we do", "OUTPUTS — What we deliver"],
    [
        ["750 facilitators", "Parenting group sessions", "85,000 caregivers reached"],
        ["Curriculum materials", "Home visits", "36 sessions/group/year"],
        ["USD 285M budget", "Health center growth monitoring", "750 facilitators deployed"],
    ],
    slide_num=6,
    note_text="Key distinction: Outputs tell us what was done. They do NOT tell us whether anything changed. This is why we need a Theory of Change."
)

# ── SLIDE 7: What are results? ──
make_content_slide(
    "What Are Results? Outcomes and Impact",
    [
        ("OUTCOME", {"bold": True, "color": MID_BLUE, "size": 22}),
        ("A change in knowledge, attitude, behavior, or condition of a person or system resulting from the intervention.", {}),
        ("", {"size": 8}),
        ("IMPACT", {"bold": True, "color": MID_BLUE, "size": 22}),
        ("Long-term, higher-order change attributable (in part) to the intervention \u2014 often at population level.", {}),
        ("", {"size": 8}),
        ("In Creciendo Juntos:", {"bold": True}),
        ("  Outcome = Caregiver changes parenting behavior", {}),
        ("  Impact   = Child achieves age-appropriate development by age 3", {}),
        ("", {"size": 8}),
        ("\u201cThe program delivers sessions (output). Sessions are supposed to change caregivers (outcome). Changed caregivers are supposed to develop children (impact). But HOW does each step happen?\u201d", {"italic": True, "color": ACCENT_ORANGE}),
    ],
    slide_num=7
)

# ── SLIDE 8: The results chain ──
make_content_slide(
    "The Results Chain: A Causal Map of Change",
    [
        ("INPUTS  \u2192  ACTIVITIES  \u2192  OUTPUTS  \u2192  OUTCOMES  \u2192  IMPACT", {"bold": True, "size": 22, "color": DARK_BLUE}),
        ("", {"size": 12}),
        ("INPUTS: 750 facilitators \u00b7 curriculum \u00b7 USD 285M", {"bold": True}),
        ("ACTIVITIES: Parenting sessions \u00b7 home visits \u00b7 health center visits", {"bold": True}),
        ("OUTPUTS: 85,000 caregivers reached \u00b7 36 sessions per group per year", {"bold": True}),
        ("OUTCOMES: Caregiver behavior change (interaction quality, nutrition practices)", {"bold": True}),
        ("IMPACT: 65% of children achieve age-appropriate development by 24 months", {"bold": True}),
    ],
    slide_num=8,
    note_text="Draw the chain on the board. Ask: 'Which of these boxes is easiest to measure? Which is hardest to attribute to the program?' — sets up the contribution analysis thread."
)

# ── SLIDE 9: Direct and distal outcomes ──
make_content_slide(
    "Direct and Distal Outcomes: Why the Distinction Matters",
    [
        ("DIRECT OUTCOME", {"bold": True, "color": MID_BLUE, "size": 22}),
        ("Change in the group directly reached by the program", {}),
        ("\u2192  Caregivers gain knowledge, change behavior", {"color": MED_GRAY}),
        ("", {"size": 8}),
        ("DISTAL OUTCOME", {"bold": True, "color": MID_BLUE, "size": 22}),
        ("Change in the group ultimately targeted", {}),
        ("\u2192  Children show improved development", {"color": MED_GRAY}),
        ("", {"size": 12}),
        ("THE INDIRECT PATHWAY", {"bold": True, "color": ACCENT_ORANGE, "size": 22}),
        ("Program  \u2192  Caregiver (direct)  \u2192  Child (distal)", {"bold": True}),
        ("", {"size": 8}),
        ("The program does NOT directly develop children\u2019s brains. It develops caregivers\u2019 capacity to develop children\u2019s brains. This is why caregiver behavioral change is the critical mechanism.", {"italic": True}),
    ],
    slide_num=9
)

# ── SLIDE 10: Monitoring and Evaluation ──
make_table_slide(
    "Monitoring and Evaluation: Different Questions, Shared Theory",
    ["", "MONITORING", "EVALUATION"],
    [
        ["Question", "Are we delivering as planned?", "Are outcomes occurring as expected?"],
        ["Focus", "Tracks outputs and early outcomes over time", "Tests whether the ToC holds \u2014 and why"],
        ["Signal", "Early warning: is the theory unfolding?", "Strength of the contribution claim"],
        ["Timing", "Continuous \u2014 during implementation", "Periodic \u2014 at key milestones"],
    ],
    slide_num=10,
    note_text="Both M&E use the same Theory of Change. Building a good ToC today is foundational to everything that follows this week."
)

# ── SLIDE 11: Exercise 1 ──
make_exercise_slide(
    "Exercise 1 \u2014 Map the Results Chain for Creciendo Juntos",
    [
        ("15 min  (10 min groups + 5 min plenary)  \u00b7  Small groups (3\u20134)", {"bold": True, "size": 16, "color": MED_GRAY}),
        ("", {"size": 8}),
        ("Using the 1-page vignette distributed at the start:", {"bold": True}),
        ("", {"size": 6}),
        ("1.  Complete the results chain template for Creciendo Juntos.", {}),
        ("    For each level, identify: (a) What is happening? (b) Who is involved? (c) Direct or distal?", {"size": 16}),
        ("", {"size": 6}),
        ("2.  Identify ONE change that was probably NOT planned by the program.", {}),
        ("", {"size": 6}),
        ("3.  Be ready to share: What is the program trying to change, and in whom?", {}),
        ("", {"size": 12}),
        ("Materials: 1-page case vignette  \u00b7  Results chain blank template", {"italic": True, "size": 16, "color": MED_GRAY}),
    ],
    slide_num=11,
    time_label="\u23f1 15 min",
    note_text="Answer key — Inputs: facilitators, curriculum, budget. Outputs: 85K caregivers, 36 sessions/yr. Direct outcomes: caregiver knowledge/behavior change. Distal outcomes: child development. Unintended: facilitator professional credentials, peer networks reduce maternal isolation."
)

# ── SLIDE 12: Debrief Ex 1 ──
make_debrief_slide(
    "Debrief: The Indirect Pathway Is the Point",
    [
        ("Most groups correctly identify outputs but struggle with the mechanism between outputs and outcomes \u2014 this gap is intentional.", {}),
        ("", {"size": 8}),
        ("Critical distinction:", {"bold": True}),
        ("The program delivers services to caregivers (direct) and expects children to benefit (distal).", {}),
        ("", {"size": 8}),
        ("What is missing from the results chain:", {"bold": True}),
        ("It tells us WHAT should happen, not WHY or HOW.", {"color": ACCENT_ORANGE}),
        ("", {"size": 12}),
        ("\u201cWhat you just built is a results chain. It is necessary \u2014 but not sufficient. A Theory of Change adds three things: causal mechanisms, causal assumptions, and external influences.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=12
)

# ── SLIDE 13: Unintended outcomes ──
make_content_slide(
    "Unintended Outcomes and Evaluation Breadth",
    [
        ("Programs regularly generate effects beyond the intended results chain.", {}),
        ("", {"size": 8}),
        ("A good Theory of Change flags plausible unintended pathways \u2014 not to manage all of them, but to avoid being surprised.", {}),
        ("", {"size": 8}),
        ("Examples from Creciendo Juntos:", {"bold": True}),
        ("  \u2022  Facilitators develop professional credentials", {}),
        ("  \u2022  Peer networks reduce maternal depression and isolation", {}),
        ("", {"size": 12}),
        ("Transition: We know WHAT should change. But WHY should it change \u2014 and HOW exactly?", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=13
)

# ── BLOCK 3 SECTION DIVIDER ──
make_section_slide("BLOCK 3", "Theory of Change", "~70 minutes  \u00b7  Slides 14\u201321", slide_num=None)

# ── SLIDE 14: What is a Theory of Change? ──
make_content_slide(
    "What Is a Theory of Change?",
    [
        ("DEFINITION", {"bold": True, "color": ACCENT_ORANGE, "size": 22}),
        ("A Theory of Change is a causal argument that explains how and why an intervention is expected to produce change.", {}),
        ("", {"size": 6}),
        ("It specifies:", {"bold": True}),
        ("  (1) The results chain    (2) The causal mechanisms", {}),
        ("  (3) The causal assumptions    (4) External influences", {}),
        ("", {"size": 12}),
        ("What makes it different from a logic model:", {"bold": True, "color": DARK_BLUE, "size": 20}),
        ("  \u2022  A logic model describes; a ToC argues", {}),
        ("  \u2022  A logic model says \u201cif we do X, then Y\u201d", {}),
        ("    A ToC says \u201cif we do X, then Y, BECAUSE OF Z, assuming A and B hold\u201d", {}),
        ("  \u2022  A ToC makes the mechanism and assumptions explicit \u2014 and therefore testable", {}),
        ("", {"size": 8}),
        ("Rogers (2008); Koleros & Mayne (2017)", {"italic": True, "size": 14, "color": MED_GRAY}),
    ],
    slide_num=14
)

# ── SLIDE 15: What is missing from the results chain? ──
make_content_slide(
    "What Is Missing from the Results Chain?",
    [
        ("Results chain says:", {"bold": True, "color": MED_GRAY}),
        ("Sessions delivered  \u2192  Caregiver behavior change  \u2192  Child development", {"bold": True}),
        ("", {"size": 12}),
        ("But does NOT explain:", {"bold": True, "color": ACCENT_ORANGE, "size": 22}),
        ("", {"size": 6}),
        ("WHY would a session change caregiver behavior?", {"bold": True}),
        ("\u2192  Causal mechanism", {"color": LIGHT_BLUE}),
        ("", {"size": 6}),
        ("WHAT must be true for that change to happen?", {"bold": True}),
        ("\u2192  Causal assumptions", {"color": LIGHT_BLUE}),
        ("", {"size": 6}),
        ("WHAT else could produce the observed outcome?", {"bold": True}),
        ("\u2192  External influences", {"color": LIGHT_BLUE}),
        ("", {"size": 10}),
        ("\u201cThe jump from sessions delivered to child development improved is not self-evident. It is a theory. And like any theory, it can be wrong.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=15
)

# ── SLIDE 16: Causal mechanisms ──
make_content_slide(
    "Causal Mechanisms: How Does Change Actually Happen?",
    [
        ("A causal mechanism is the process through which an output produces an outcome.", {"bold": True}),
        ("", {"size": 6}),
        ("Key distinctions:", {"bold": True, "color": DARK_BLUE}),
        ("  \u2022  Mechanisms are NOT activities \u2014 activities produce outputs; mechanisms explain how outputs produce outcomes", {}),
        ("  \u2022  In theory-based evaluation, we test mechanisms directly \u2014 not just outcomes", {}),
        ("", {"size": 12}),
        ("EXAMPLE: CRECIENDO JUNTOS", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("  Activity: 36 parenting group sessions delivered", {}),
        ("  Mechanism: Sessions increase caregiver CAPABILITY (knowledge), create OPPORTUNITY (peer networks), build MOTIVATION (parental identity)", {}),
        ("  Outcome: Caregivers change daily interaction behaviors", {}),
        ("  Distal outcome: Children experience improved stimulation \u2192 developmental gains", {}),
        ("", {"size": 8}),
        ("\u201cThe mechanism is what sessions are supposed to trigger in caregivers. If the mechanism does not fire, sessions may be delivered but nothing changes.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=16
)

# ── SLIDE 17: COM-B ──
make_table_slide(
    "COM-B: Capability, Opportunity, Motivation \u2192 Behavior",
    ["Component", "Meaning", "In Creciendo Juntos"],
    [
        ["CAPABILITY (C)", "What a person knows and can do", "Caregiver learns developmental milestones, responsive play"],
        ["OPPORTUNITY (O)", "What the environment enables or constrains", "Peer group forms, materials distributed, health visits accessible"],
        ["MOTIVATION (M)", "What drives action", "Child responds positively; new identity: \"I am my child's first teacher\""],
        ["\u2192 BEHAVIOR", "Observable action that produces the outcome", "Daily responsive interaction: 15 \u2192 45 min; harsh discipline replaced"],
    ],
    slide_num=17,
    note_text="COM-B is an analytical lens, not a checklist. Its value is forcing the question: what exactly needs to shift in the person for new behavior to become possible?"
)

# ── SLIDE 18: COM-B in CJ ──
make_content_slide(
    "COM-B in Creciendo Juntos: Tracing the Pathway",
    [
        ("CJ Sessions  \u2192  [C] + [O] + [M]  \u2192  Behavior change  \u2192  Child development", {"bold": True, "color": DARK_BLUE, "size": 20}),
        ("", {"size": 10}),
        ("CAPABILITY", {"bold": True, "color": LIGHT_BLUE}),
        ("Knowledge of child development stages  \u00b7  Skill in responsive interaction  \u00b7  Self-efficacy", {}),
        ("", {"size": 6}),
        ("OPPORTUNITY", {"bold": True, "color": LIGHT_BLUE}),
        ("Peer support group  \u00b7  Learning materials  \u00b7  Health center access", {}),
        ("", {"size": 6}),
        ("MOTIVATION", {"bold": True, "color": LIGHT_BLUE}),
        ("Emotional reward from child\u2019s positive response  \u00b7  Peer norms  \u00b7  New parental identity", {}),
        ("", {"size": 6}),
        ("BEHAVIOR", {"bold": True, "color": ACCENT_ORANGE}),
        ("Responsive interaction \u2191  \u00b7  Positive discipline  \u00b7  Enriched home environment", {}),
        ("", {"size": 10}),
        ("\u201cWe traced the mechanism for one actor \u2014 the caregiver. In the exercise, you will go deeper.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=18
)

# ── SLIDE 19: Exercise 2 ──
make_exercise_slide(
    "Exercise 2 \u2014 Build the Causal Pathway for Creciendo Juntos",
    [
        ("20 min  (15 min groups + 5 min plenary)  \u00b7  Pairs or small groups", {"bold": True, "size": 16, "color": MED_GRAY}),
        ("", {"size": 8}),
        ("Choose ONE reached group: (a) caregivers, or (b) community health workers.", {"bold": True}),
        ("", {"size": 6}),
        ("For your chosen group, trace the complete causal pathway:", {}),
        ("  (a) CAPABILITY: What must this person know or be able to do?", {}),
        ("  (b) OPPORTUNITY: What environmental conditions must be present?", {}),
        ("  (c) MOTIVATION: What must shift in their sense of purpose or habit?", {}),
        ("  (d) BEHAVIOR: What specific observable action should change?", {}),
        ("  (e) OUTCOME: What change does this behavior produce in the target group?", {}),
        ("", {"size": 8}),
        ("Add your pathway to the results chain template from Exercise 1.", {"bold": True}),
        ("", {"size": 8}),
        ("Materials: Results chain template from Ex. 1  \u00b7  COM-B lens card", {"italic": True, "size": 16, "color": MED_GRAY}),
    ],
    slide_num=19,
    time_label="\u23f1 20 min",
    note_text="Answer key (caregivers): Capability = milestones knowledge, responsive skills. Opportunity = peer group, materials, health services. Motivation = child response, peer norms, identity. Behavior = 45 min interaction, language-rich, no harsh discipline. Outcome = developmental gains by 24 months."
)

# ── SLIDE 20: Debrief Ex 2 ──
make_debrief_slide(
    "Debrief: Connecting Outputs \u2192 Behaviors \u2192 Outcomes",
    [
        ("The behavioral mechanism is what sessions are supposed to trigger. Without it, outputs do not produce outcomes.", {}),
        ("", {"size": 8}),
        ("Different groups trace different pathways \u2014 both matter.", {}),
        ("", {"size": 8}),
        ("The pathway you traced is a hypothesis. Each link depends on something being true.", {"bold": True}),
        ("", {"size": 12}),
        ("\u201cYou have built the core of the Theory of Change. The next question is: what are we assuming? What must be true for each link to hold?\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=20
)

# ── SLIDE 21: Full ToC diagram ──
make_content_slide(
    "The Creciendo Juntos ToC: The Full Causal Picture",
    [
        ("Four layers of the Theory of Change:", {"bold": True, "color": DARK_BLUE, "size": 22}),
        ("", {"size": 10}),
        ("\u2588\u2588  Layer 1 \u2014 Results chain (Exercise 1)", {"bold": True, "color": DARK_GRAY}),
        ("    INPUTS \u2192 ACTIVITIES \u2192 OUTPUTS \u2192 OUTCOMES \u2192 IMPACT", {"size": 16}),
        ("", {"size": 6}),
        ("\u2588\u2588  Layer 2 \u2014 Behavioral mechanism / COM-B (Exercise 2)", {"bold": True, "color": LIGHT_BLUE}),
        ("    Capability + Opportunity + Motivation \u2192 Behavior", {"size": 16}),
        ("", {"size": 6}),
        ("\u2588\u2588  Layer 3 \u2014 Causal assumptions, classified by level (Exercise 3)", {"bold": True, "color": ACCENT_ORANGE}),
        ("    Reach \u00b7 Capacity \u00b7 Behavior \u00b7 Outcome", {"size": 16}),
        ("", {"size": 6}),
        ("\u2588\u2588  Layer 4 \u2014 External influences (Exercise 4)", {"bold": True, "color": MED_GRAY}),
        ("    Contextual factors \u00b7 Rival explanations", {"size": 16}),
        ("", {"size": 10}),
        ("\u201cWhat surprises you? What would you add or question?\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=21,
    note_text="Ask: 'What surprises you? What would you add or question?' Take 2\u20133 responses. Keep brief \u2014 Block 4 will go deep on assumptions."
)

# ── BLOCK 4 SECTION DIVIDER ──
make_section_slide("BLOCK 4", "Assumptions and External Influences", "~65 minutes  \u00b7  Slides 22\u201329", slide_num=None)

# ── SLIDE 22: Causal assumptions ──
make_content_slide(
    "What Are Causal Assumptions?",
    [
        ("DEFINITION", {"bold": True, "color": ACCENT_ORANGE, "size": 22}),
        ("A causal assumption is a condition that must hold for one level of the results chain to produce the next.", {}),
        ("If an assumption is violated, the expected causal link breaks \u2014 even if the program is implemented correctly.", {}),
        ("", {"size": 10}),
        ("EXAMPLE FROM CRECIENDO JUNTOS", {"bold": True, "color": DARK_BLUE, "size": 20}),
        ("  Link: Parenting sessions delivered \u2192 Caregivers change behavior", {}),
        ("  Assumption: \u201cAttendance is sufficient to build capability and motivation", {}),
        ("  (caregivers attend \u2265 70% of sessions)\u201d", {}),
        ("", {"size": 6}),
        ("  If violated: Sessions delivered but caregivers attend sporadically \u2192 capability does not build \u2192 behavior does not change \u2192 outcomes not reached", {"color": ACCENT_ORANGE}),
        ("", {"size": 10}),
        ("\u201cThis is why monitoring attendance is not optional. If this assumption breaks, everything downstream breaks with it.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=22
)

# ── SLIDE 23: Types of assumptions ──
make_table_slide(
    "Types of Assumptions: Classifying by Results Chain Level",
    ["Assumption Type", "Question It Answers", "CJ Example"],
    [
        ["REACH", "Can the program reach the intended population?", "Facilitators can be recruited and retained in target communities"],
        ["CAPACITY", "Does the intervention build the required COM-B shift?", "Sessions delivered with sufficient quality to build caregiver capability"],
        ["BEHAVIOR", "Does changed COM-B produce behavior change?", "Knowledge and motivation translate to daily behavior change (not just during sessions)"],
        ["OUTCOME", "Does behavior change produce the expected outcome?", "Changed caregiver behavior sustained 18+ months produces developmental gains"],
    ],
    slide_num=23
)

# ── SLIDE 24: Critical assumptions ──
make_content_slide(
    "Critical Assumptions: Which Links Are Most Fragile?",
    [
        ("A critical assumption is one that:", {"bold": True, "color": DARK_BLUE, "size": 20}),
        ("  1.  Is ESSENTIAL \u2014 if it breaks, outcomes cannot be reached regardless of delivery quality", {}),
        ("  2.  Is UNCERTAIN \u2014 we have limited evidence that it holds in this context", {}),
        ("  3.  Is TESTABLE \u2014 we can design monitoring or evaluation to check it", {}),
        ("", {"size": 12}),
        ("TWO CANDIDATE CRITICAL ASSUMPTIONS", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("", {"size": 6}),
        ("Assumption B1 (BEHAVIOR):", {"bold": True}),
        ("\u201cCaregivers who increase knowledge in sessions also change their daily parenting behaviors at home, not just during observed sessions.\u201d", {"italic": True}),
        ("CRITICAL because behavioral generalization is uncertain and rarely measured.", {"color": MED_GRAY, "size": 16}),
        ("", {"size": 8}),
        ("Assumption O1 (OUTCOME):", {"bold": True}),
        ("\u201cChanged parenting behavior, sustained for 18 months, is sufficient to produce measurable developmental gains in children by age 3.\u201d", {"italic": True}),
        ("CRITICAL because the dose-response relationship is unclear for multi-risk households.", {"color": MED_GRAY, "size": 16}),
    ],
    slide_num=24
)

# ── SLIDE 25: Exercise 3 ──
make_exercise_slide(
    "Exercise 3 \u2014 Identify and Classify Assumptions in the ToC",
    [
        ("20 min  (15 min groups + 5 min plenary)  \u00b7  Small groups (3\u20134)", {"bold": True, "size": 16, "color": MED_GRAY}),
        ("", {"size": 8}),
        ("Using the Creciendo Juntos ToC from Slide 21:", {"bold": True}),
        ("", {"size": 6}),
        ("1.  Identify at least ONE assumption per level: reach, capacity, behavior, outcome.", {}),
        ("", {"size": 4}),
        ("2.  Write each assumption as a complete sentence:", {}),
        ("    \u201cIF [condition], THEN the link holds.\u201d", {"italic": True}),
        ("", {"size": 4}),
        ("3.  Mark which TWO assumptions you consider most critical for this program\u2019s success.", {}),
        ("", {"size": 4}),
        ("4.  Be ready to explain: Why did you choose these two? What evidence would test them?", {}),
        ("", {"size": 8}),
        ("Add your assumptions to the ToC template \u2014 they should annotate each causal arrow.", {"bold": True}),
        ("", {"size": 8}),
        ("Materials: ToC template (annotated from Ex. 1 & 2)  \u00b7  Assumption classification card", {"italic": True, "size": 16, "color": MED_GRAY}),
    ],
    slide_num=25,
    time_label="\u23f1 20 min",
    note_text="The assumptions identified are exactly the questions a contribution analysis must test. In Session 5, we will ask: which assumptions held? How strong is our contribution claim?"
)

# ── SLIDE 26: Debrief Ex 3 ──
make_debrief_slide(
    "Debrief: Which Assumptions Are Most Fragile?",
    [
        ("Most groups identify Behavior assumptions as most critical \u2014 good: this is where most ECD programs fail.", {}),
        ("", {"size": 8}),
        ("An assumption is only useful if it is testable.", {"bold": True}),
        ("", {"size": 6}),
        ("Critical assumptions become monitoring indicators AND evaluation questions.", {}),
        ("", {"size": 12}),
        ("CONNECTION TO EVALUATION QUESTIONS", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("", {"size": 6}),
        ("\u201cIf the critical assumption is \u2018caregivers change daily behavior,\u2019 the evaluation question is:", {"italic": True, "color": MID_BLUE}),
        ("Did caregivers actually change their daily parenting behavior \u2014 and if so, how much, for whom, and under what conditions?", {"italic": True, "color": MID_BLUE}),
        ("That is what Session 5 will help you answer.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=26
)

# ── SLIDE 27: External influences ──
make_content_slide(
    "External Influences: What Is Outside the Program\u2019s Control?",
    [
        ("CONTEXTUAL FACTORS", {"bold": True, "color": MID_BLUE, "size": 22}),
        ("Structural conditions that shape whether the mechanism can function:", {}),
        ("  economy, political stability, health system quality", {"color": MED_GRAY}),
        ("", {"size": 10}),
        ("RIVAL EXPLANATIONS", {"bold": True, "color": MID_BLUE, "size": 22}),
        ("Alternative causes of the observed outcome that would have occurred without the program:", {}),
        ("  natural development, another program, a social trend", {"color": MED_GRAY}),
        ("", {"size": 12}),
        ("WHY BOTH MATTER FOR EVALUATION", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("Contextual factors can violate assumptions \u2014 e.g., severe food insecurity may mean that even improved parenting behavior does not translate to developmental gains.", {}),
        ("", {"size": 6}),
        ("Rival explanations: if children improve even in control communities, the program\u2019s contribution is less clear. Evaluators must rule out or account for them.", {}),
    ],
    slide_num=27
)

# ── SLIDE 28: Exercise 4 ──
make_exercise_slide(
    "Exercise 4 \u2014 External Influences Scan (Full-Group Discussion)",
    [
        ("10 min  \u00b7  Facilitated group discussion", {"bold": True, "size": 16, "color": MED_GRAY}),
        ("", {"size": 8}),
        ("Q1: What are the 2\u20133 most important external factors that could affect whether Creciendo Juntos achieves its outcomes?", {"bold": True}),
        ("    (Facilitator records on board)", {"size": 16, "color": MED_GRAY}),
        ("", {"size": 6}),
        ("Q2: Which of these could be a rival explanation for observed improvements?", {"bold": True}),
        ("    Which could make outcomes improve even if the program had no effect?", {"size": 16}),
        ("", {"size": 6}),
        ("Q3: Which critical assumption from Exercise 3 is most vulnerable to these external factors?", {"bold": True}),
        ("", {"size": 12}),
        ("Expected contributions:", {"bold": True, "color": DARK_BLUE}),
        ("  \u2022  Economic conditions (recession \u2192 food insecurity)", {}),
        ("  \u2022  Political instability (budget cut, municipal support withdrawn)", {}),
        ("  \u2022  Other social programs (CCTs improving household conditions)", {}),
        ("  \u2022  Natural disasters / pandemics", {}),
    ],
    slide_num=28,
    time_label="\u23f1 10 min",
    note_text="These factors go on the edge of the ToC diagram. In evaluation design, we call these threats to internal validity. We return to them in Sessions 3\u20135."
)

# ── SLIDE 29: Complete ToC ──
make_content_slide(
    "The Complete Theory of Change: Four Layers",
    [
        ("", {"size": 4}),
        ("\u2588\u2588  LAYER 1 \u2014 Results chain (Exercise 1)", {"bold": True, "color": DARK_GRAY, "size": 20}),
        ("    Inputs \u2192 Activities \u2192 Outputs \u2192 Outcomes \u2192 Impact", {"size": 16}),
        ("", {"size": 6}),
        ("\u2588\u2588  LAYER 2 \u2014 Behavioral mechanism / COM-B (Exercise 2)", {"bold": True, "color": LIGHT_BLUE, "size": 20}),
        ("    C + O + M \u2192 Behavior change", {"size": 16}),
        ("", {"size": 6}),
        ("\u2588\u2588  LAYER 3 \u2014 Causal assumptions (Exercise 3)", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("    Classified: Reach \u00b7 Capacity \u00b7 Behavior \u00b7 Outcome", {"size": 16}),
        ("", {"size": 6}),
        ("\u2588\u2588  LAYER 4 \u2014 External influences (Exercise 4)", {"bold": True, "color": MED_GRAY, "size": 20}),
        ("    Contextual factors + Rival explanations", {"size": 16}),
        ("", {"size": 12}),
        ("\u201cThis is a Theory of Change. It is not a diagram \u2014 it is an argument. You can now use it to design evaluation questions, select indicators, and assess your contribution claim.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=29
)

# ── BLOCK 5 SECTION DIVIDER ──
make_section_slide("BLOCK 5", "M&E and Bridge to Session 2", "~30 minutes  \u00b7  Slides 30\u201341", slide_num=None)

# ── SLIDE 30: ToC and evaluation questions ──
make_table_slide(
    "ToC and Evaluation Questions: Making the Theory Testable",
    ["Critical Assumption (from Exercise 3)", "Evaluation Question It Generates"],
    [
        [
            "Behavior: Caregivers change daily practices, not just session behavior",
            "Did caregivers change their daily interaction behavior at home? For whom? To what extent?"
        ],
        [
            "Outcome: Changed behavior, sustained 18 months, produces developmental gains",
            "Is there a relationship between intensity of caregiver behavior change and child outcomes?"
        ],
        [
            "Reach: Facilitators can be retained for program duration",
            "What is the facilitator turnover rate and how does it affect session continuity and trust?"
        ],
    ],
    slide_num=30,
    note_text="You did not just build a diagram. You built a set of evaluation questions. Evaluation design starts here \u2014 from the assumptions."
)

# ── SLIDE 31: Monitoring ──
make_content_slide(
    "Monitoring: Watching the ToC Unfold in Real Time",
    [
        ("Monitoring observes: reach, outputs, and early behavioral signals.", {}),
        ("", {"size": 8}),
        ("Monitoring does not test the ToC \u2014 it signals whether conditions for the ToC to function are present.", {"bold": True}),
        ("", {"size": 8}),
        ("Example:", {"bold": True, "color": DARK_BLUE}),
        ("If monitoring shows low attendance \u2192 reach assumption may be breaking \u2192 evaluator can investigate before outcomes are affected.", {}),
        ("", {"size": 12}),
        ("CONNECTION TO SESSION 2", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("\u201cWhat you monitor is determined by your ToC. In Session 2, we will take this results chain and ask: what indicator would tell us whether each level is happening? That is the Results Matrix.\u201d", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=31
)

# ── SLIDE 32: Exercise 5 ──
make_exercise_slide(
    "Exercise 5 \u2014 What Would You Monitor First? (Optional)",
    [
        ("8 min  (5 min pairs + 3 min share-out)  \u00b7  Can be assigned as homework", {"bold": True, "size": 16, "color": MED_GRAY}),
        ("", {"size": 8}),
        ("Looking at the completed Creciendo Juntos ToC:", {"bold": True}),
        ("", {"size": 6}),
        ("1.  Choose the THREE variables you would monitor most closely in the first 6 months.", {}),
        ("", {"size": 4}),
        ("2.  For each: (a) What would you measure? (b) Who would you observe or survey? (c) When?", {}),
        ("", {"size": 4}),
        ("3.  Explain: Why these three? What would each tell you about whether the ToC is unfolding?", {}),
        ("", {"size": 10}),
        ("Note: These are proto-indicators. In Session 2, we will formalize them into a Results Matrix.", {"italic": True, "color": MID_BLUE}),
    ],
    slide_num=32,
    time_label="\u23f1 8 min",
    note_text="Expected answers: Facilitator turnover rate (reach) \u2014 monthly HR records. Session attendance (capacity) \u2014 per group per session. Observed caregiver interaction (behavior) \u2014 structured observation at Month 3 and 6."
)

# ── SLIDE 33: Evaluation at outcome level ──
make_content_slide(
    "Evaluation: Testing the ToC at Outcome Level",
    [
        ("Evaluation asks:", {"bold": True, "color": DARK_BLUE, "size": 20}),
        ("  Did expected outcomes occur?", {}),
        ("  Did the theory\u2019s mechanisms function as hypothesized?", {}),
        ("", {"size": 10}),
        ("Contribution analysis (Session 5):", {"bold": True, "color": DARK_BLUE, "size": 20}),
        ("\u201cGiven what we observe, how strong is our claim that the program contributed to these outcomes?\u201d", {"italic": True}),
        ("", {"size": 8}),
        ("The ToC built today is the foundation of the contribution claim \u2014 every assumption identified is a piece of evidence to test.", {"bold": True}),
        ("", {"size": 12}),
        ("THREAD TO SESSIONS 3\u20135", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("Sessions 3 and 4: What qualitative and quantitative evidence would let you test whether each mechanism and assumption held?", {}),
        ("Session 5: Brings all evidence into a contribution claim.", {}),
        ("Today\u2019s ToC is the map.", {"bold": True, "color": MID_BLUE}),
    ],
    slide_num=33
)

# ── SLIDE 34: From ToC to Results Matrix ──
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, TRANSITION_GREEN)

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = ACCENT_GREEN
bar.line.fill.background()

tb = add_text_box(s, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9))
set_text(tb.text_frame, "From ToC to Results Matrix: What Comes Next", size=28, bold=True, color=DARK_BLUE)

# Left column: today you built
lb = add_text_box(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(3))
lb.text_frame.word_wrap = True
set_text(lb.text_frame, "TODAY YOU BUILT:", size=20, bold=True, color=ACCENT_GREEN)
for t in [
    "\u2713  A results chain (Exercise 1)",
    "\u2713  A causal mechanism layer \u2014 COM-B (Exercise 2)",
    "\u2713  A classified assumption list (Exercise 3)",
    "\u2713  An external influence map (Exercise 4)",
]:
    add_paragraph(lb.text_frame, t, size=17, color=DARK_GRAY)

# Right column: Session 2 preview
rb = add_text_box(s, Inches(6.8), Inches(1.4), Inches(5.8), Inches(3))
rb.text_frame.word_wrap = True
set_text(rb.text_frame, "SESSION 2 WILL ASK:", size=20, bold=True, color=DARK_BLUE)
add_paragraph(rb.text_frame, "\u201cWhat indicator would tell us whether each level is happening?\u201d", size=17, italic=True, color=MID_BLUE)
add_paragraph(rb.text_frame, "\u201cWhat source? What frequency? What baseline?\u201d", size=17, italic=True, color=MID_BLUE)

# Preview table
headers = ["Element", "Content"]
rows_data = [
    ["Results chain level", "Caregiver behavior change"],
    ["Indicator", "% of caregivers observed engaging in \u2265 3 responsive interaction activities per hour"],
    ["Type", "Behavioral (mechanism) indicator"],
    ["Baseline", "15% of caregivers at enrollment"],
    ["Target", "65% at Month 9"],
    ["Source", "Structured observation at home visits (Months 3, 6, 9)"],
]
tbl = s.shapes.add_table(len(rows_data) + 1, 2, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.5)).table
tbl.columns[0].width = Inches(3)
tbl.columns[1].width = Inches(8.5)
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = ACCENT_GREEN
for i, row in enumerate(rows_data):
    for j, val in enumerate(row):
        cell = tbl.cell(i + 1, j)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE

sn = add_text_box(s, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
set_text(sn.text_frame, "34", size=12, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)

# ── CLOSURE SECTION ──
make_section_slide("CLOSURE", "Takeaways \u00b7 Survey \u00b7 Pre-reading", "~15 minutes  \u00b7  Slides 35\u201344", slide_num=None)

# ── SLIDE 35: Key takeaways ──
make_content_slide(
    "Key Takeaways: Returning to the Learning Objectives",
    [
        ("\u2713  Frame complex interventions", {"bold": True, "color": ACCENT_GREEN, "size": 20}),
        ("   Creciendo Juntos reaches caregivers to change children; the indirect pathway is the mechanism", {"size": 16}),
        ("", {"size": 6}),
        ("\u2713  Explain theory-based evaluation", {"bold": True, "color": ACCENT_GREEN, "size": 20}),
        ("   The results chain alone is insufficient; mechanisms, assumptions, and external influences are essential", {"size": 16}),
        ("", {"size": 6}),
        ("\u2713  Describe a Theory of Change", {"bold": True, "color": ACCENT_GREEN, "size": 20}),
        ("   Four layers: results chain, mechanisms, assumptions, external influences", {"size": 16}),
        ("", {"size": 6}),
        ("\u2713  Distinguish M&E", {"bold": True, "color": ACCENT_GREEN, "size": 20}),
        ("   Monitoring watches the ToC unfold; evaluation tests whether it held", {"size": 16}),
        ("", {"size": 6}),
        ("\u2713  Formulate evaluation questions", {"bold": True, "color": ACCENT_GREEN, "size": 20}),
        ("   Each critical assumption generates a specific evaluative question", {"size": 16}),
    ],
    slide_num=35,
    note_text="Closing question for participants: 'Think of a program you work with. What is the most critical assumption in its Theory of Change?'"
)

# ── SLIDE 36: Pre-reading ──
make_content_slide(
    "Pre-reading for Session 2",
    [
        ("BEFORE SESSION 2 \u2014 READ:", {"bold": True, "color": ACCENT_ORANGE, "size": 24}),
        ("", {"size": 12}),
        ("Creciendo Juntos: Narrative Part 1, Pages 1\u20134", {"bold": True, "size": 20}),
        ("(Program design and implementation structure)", {"color": MED_GRAY}),
        ("", {"size": 12}),
        ("Focus on:", {"bold": True}),
        ("  \u2022  How are the four components coordinated?", {}),
        ("  \u2022  What roles do facilitators, health staff, and coordinators play?", {}),
        ("", {"size": 12}),
        ("In Session 2, we will use this program design to populate a full Results Matrix \u2014 one indicator set per results chain level.", {"italic": True}),
        ("", {"size": 12}),
        ("Estimated reading time: 15\u201320 minutes.", {"color": MED_GRAY}),
    ],
    slide_num=36
)

# ── SLIDE 37: Closing survey ──
make_content_slide(
    "Closing Survey",
    [
        ("Standard course satisfaction items (per OVE format)", {"color": MED_GRAY}),
        ("", {"size": 12}),
        ("EVALUATIVE QUESTIONS:", {"bold": True, "color": DARK_BLUE, "size": 22}),
        ("", {"size": 8}),
        ("Q1: In one sentence: what is the causal mechanism in the Creciendo Juntos Theory of Change?", {"bold": True}),
        ("", {"size": 8}),
        ("Q2: What is the most critical assumption in a Theory of Change you currently work with?", {"bold": True}),
        ("", {"size": 12}),
        ("These are formative assessment items that feed directly into the EVALAC pre/post evaluation.", {"italic": True, "color": MED_GRAY, "size": 16}),
    ],
    slide_num=37
)

# ── SLIDES 38-44: Reference slides ──
ref_slides = [
    (38, "Reference: Glossary of Key Concepts", [
        ("Results chain: The logical sequence linking inputs \u2192 activities \u2192 outputs \u2192 outcomes \u2192 impact", {}),
        ("Theory of Change: A causal argument specifying the results chain, causal mechanisms, assumptions, and external influences", {}),
        ("Causal mechanism: The process through which an output produces an outcome", {}),
        ("COM-B: Capability, Opportunity, Motivation \u2192 Behavior \u2014 a framework for analyzing behavioral mechanisms", {}),
        ("Causal assumption: A condition that must hold for one level of the results chain to produce the next", {}),
        ("Critical assumption: An assumption that is essential, uncertain, and testable", {}),
        ("Direct outcome: Change in the group directly reached by the program", {}),
        ("Distal outcome: Change in the group ultimately targeted (through indirect pathway)", {}),
        ("Monitoring: Continuous tracking of reach, outputs, and early behavioral signals", {}),
        ("Evaluation: Periodic assessment of whether the ToC held \u2014 and why", {}),
        ("Contribution analysis: Assessment of the strength of the claim that the program contributed to observed outcomes", {}),
    ]),
    (39, "Reference: Creciendo Juntos \u2014 Full Annotated ToC", [
        ("LAYER 1 \u2014 RESULTS CHAIN", {"bold": True, "color": DARK_GRAY, "size": 20}),
        ("Inputs (facilitators, curriculum, USD 285M) \u2192 Activities (sessions, home visits, health monitoring) \u2192 Outputs (85K caregivers, 36 sessions/yr) \u2192 Outcomes (caregiver behavior change) \u2192 Impact (child development by age 3)", {"size": 15}),
        ("", {"size": 8}),
        ("LAYER 2 \u2014 BEHAVIORAL MECHANISM (COM-B)", {"bold": True, "color": LIGHT_BLUE, "size": 20}),
        ("C: Developmental knowledge + responsive skills + self-efficacy", {"size": 15}),
        ("O: Peer group + materials + health access", {"size": 15}),
        ("M: Child response + peer norms + parental identity", {"size": 15}),
        ("\u2192 B: Responsive interaction \u2191, positive discipline, enriched environment", {"size": 15}),
        ("", {"size": 8}),
        ("LAYER 3 \u2014 ASSUMPTIONS", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
        ("Reach: facilitators recruited/retained. Capacity: quality delivery. Behavior: daily transfer. Outcome: sustained 18+ mo.", {"size": 15}),
        ("", {"size": 8}),
        ("LAYER 4 \u2014 EXTERNAL INFLUENCES", {"bold": True, "color": MED_GRAY, "size": 20}),
        ("Contextual: food security, political stability, health system. Rivals: CCTs, natural development, other programs.", {"size": 15}),
    ]),
    (40, "Reference: Results Matrix Preview \u2014 One Completed Row", [
        ("This row will be expanded into a full Results Matrix in Session 2.", {"italic": True, "color": MID_BLUE}),
        ("", {"size": 8}),
        ("Results chain level:  Caregiver behavior change", {"bold": True}),
        ("Indicator:  % of caregivers observed engaging in \u2265 3 responsive interaction activities per hour during structured home visit", {}),
        ("Type:  Behavioral (mechanism) indicator", {}),
        ("Baseline:  15% of caregivers at enrollment", {}),
        ("Target:  65% at Month 9", {}),
        ("Source:  Structured observation at home visits (Months 3, 6, 9)", {}),
    ]),
    (41, "Reference: COM-B Lens Card", [
        ("CAPABILITY (C) \u2014 What a person knows and can do", {"bold": True, "color": LIGHT_BLUE}),
        ("  Physical capability: skills, strength, stamina", {}),
        ("  Psychological capability: knowledge, psychological skills, memory, decision-making", {}),
        ("", {"size": 8}),
        ("OPPORTUNITY (O) \u2014 What the environment enables or constrains", {"bold": True, "color": LIGHT_BLUE}),
        ("  Physical opportunity: time, resources, locations, cues", {}),
        ("  Social opportunity: cultural norms, social influences, interpersonal relationships", {}),
        ("", {"size": 8}),
        ("MOTIVATION (M) \u2014 What drives action", {"bold": True, "color": LIGHT_BLUE}),
        ("  Reflective motivation: plans, evaluations, identity, beliefs", {}),
        ("  Automatic motivation: emotions, desires, impulses, habits", {}),
        ("", {"size": 8}),
        ("\u2192 BEHAVIOR \u2014 The observable action that produces the intended outcome", {"bold": True, "color": ACCENT_ORANGE}),
    ]),
    (42, "Reference: Assumption Classification Card", [
        ("REACH assumptions", {"bold": True, "color": MID_BLUE, "size": 20}),
        ("Can the program reach the intended population? Are target groups accessible and willing?", {}),
        ("", {"size": 8}),
        ("CAPACITY assumptions", {"bold": True, "color": MID_BLUE, "size": 20}),
        ("Does the intervention build the required COM-B shift? Is delivery quality sufficient?", {}),
        ("", {"size": 8}),
        ("BEHAVIOR assumptions", {"bold": True, "color": MID_BLUE, "size": 20}),
        ("Does changed COM-B actually produce sustained behavior change outside the program context?", {}),
        ("", {"size": 8}),
        ("OUTCOME assumptions", {"bold": True, "color": MID_BLUE, "size": 20}),
        ("Does behavior change produce the expected developmental, social, or economic outcome?", {}),
        ("", {"size": 8}),
        ("For each: Is it ESSENTIAL? Is it UNCERTAIN? Is it TESTABLE?", {"bold": True, "color": ACCENT_ORANGE}),
    ]),
    (43, "Reference: Key Readings", [
        ("Rogers, P. (2008). Using programme theory to evaluate complicated and complex aspects of interventions. Evaluation, 14(1), 29\u201348.", {}),
        ("", {"size": 8}),
        ("Koleros, A. & Mayne, J. (2017). Using actor-based theories of change to conduct robust evaluation in complex settings. Evaluation, 25(1), 82\u2013101.", {}),
        ("", {"size": 8}),
        ("Mayne, J. (2012). Contribution analysis: Coming of age? Evaluation, 18(3), 270\u2013280.", {}),
        ("", {"size": 8}),
        ("Mayne, J. (2019). Revisiting contribution analysis. Canadian Journal of Program Evaluation, 34(2), 171\u2013191.", {}),
        ("", {"size": 8}),
        ("Michie, S., van Stralen, M.M. & West, R. (2011). The behaviour change wheel: A new method for characterising and designing behaviour change interventions. Implementation Science, 6(42).", {}),
        ("", {"size": 8}),
        ("ILAC Brief 16: Contribution Analysis. Mayne, J. (2008).", {}),
    ]),
    (44, "EVALAC I \u2014 Course References and Acknowledgments", [
        ("EVALAC I: Introduction to M&E and Theory-Based Evaluation", {"bold": True, "size": 22, "color": DARK_BLUE}),
        ("", {"size": 12}),
        ("Organized by:", {"bold": True}),
        ("  CLEAR-LAB  \u00b7  IDB  \u00b7  OVE", {}),
        ("", {"size": 8}),
        ("Case study:", {"bold": True}),
        ("  Creciendo Juntos \u2014 Integrated Early Childhood Development", {}),
        ("  (Fictional composite case based on real ECD programs in Latin America and the Caribbean)", {"size": 16, "color": MED_GRAY}),
        ("", {"size": 12}),
        ("Session 1 designed and facilitated by the EVALAC teaching team.", {}),
        ("", {"size": 8}),
        ("For questions and feedback: evalac@clearlab.org", {"italic": True, "color": MID_BLUE}),
    ]),
]

for snum, title, items in ref_slides:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, REF_BG)

    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = MED_GRAY
    bar.line.fill.background()

    tb = add_text_box(s, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8))
    set_text(tb.text_frame, title, size=24, bold=True, color=DARK_BLUE)

    cb = add_text_box(s, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    cb.text_frame.word_wrap = True
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        if first:
            p = cb.text_frame.paragraphs[0]
            first = False
        else:
            p = cb.text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(opts.get("size", 16))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK_GRAY)
        p.font.italic = opts.get("italic", False)
        p.font.name = "Calibri"
        p.space_before = Pt(4)
        p.space_after = Pt(3)

    sn = add_text_box(s, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
    set_text(sn.text_frame, str(snum), size=12, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════

out_path = os.path.join(os.path.dirname(__file__), "S1_Final_v01.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
