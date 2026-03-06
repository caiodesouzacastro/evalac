"""
Generate Session 1 PPTX v02 — EVALAC I: Introduction to M&E and Theory of Change.
Incorporates BID feedback: revised cheat sheet references, narrative handout integration,
exercise answer keys, COM-B behavioral lens, and pedagogical bridges.
44 slides, 4.5 hours, 5 exercises.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
ACCENT_ORANGE = RGBColor(0xE8, 0x6C, 0x00)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MED_GRAY = RGBColor(0x75, 0x75, 0x75)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
EXERCISE_BG = RGBColor(0xE3, 0xF2, 0xFD)
EXERCISE_ACCENT = RGBColor(0x15, 0x65, 0xC0)
DEBRIEF_BG = RGBColor(0xFE, 0xF3, 0xE0)
TRANSITION_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
REF_BG = RGBColor(0xEC, 0xEF, 0xF1)
QUOTE_BG = RGBColor(0xF3, 0xE5, 0xF5)


# ── Helpers ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, color, x=0, y=0, w=None, h=None):
    w = w or prs.slide_width
    h = h or prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT,
                   space_before=Pt(6), space_after=Pt(4), font_name="Calibri", italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_run(p, text, size=18, bold=False, color=DARK_GRAY, italic=False, font_name="Calibri"):
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = font_name
    return run


def slide_number(slide, num):
    sn = add_text_box(slide, Inches(12.2), Inches(7), Inches(1), Inches(0.4))
    set_text(sn.text_frame, str(num), size=12, color=MED_GRAY, alignment=PP_ALIGN.RIGHT)


def top_bar(slide, color=MID_BLUE, height=0.06):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def bottom_accent(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()


def make_title_slide(title, subtitle="", num=None, bg=DARK_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg)
    bottom_accent(slide)
    tb = add_text_box(slide, Inches(1), Inches(2), Inches(11), Inches(2.5))
    set_text(tb.text_frame, title, size=36, bold=True, color=WHITE)
    if subtitle:
        add_paragraph(tb.text_frame, subtitle, size=20, color=WHITE, space_before=Pt(16))
    if num:
        slide_number(slide, num)
    return slide


def make_section_slide(block_label, block_title, time_info="", num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MID_BLUE)
    tb = add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(3))
    p = tb.text_frame.paragraphs[0]
    add_run(p, block_label, size=16, bold=True, color=RGBColor(0xBB, 0xDE, 0xFB))
    add_paragraph(tb.text_frame, block_title, size=36, bold=True, color=WHITE, space_before=Pt(12))
    if time_info:
        add_paragraph(tb.text_frame, time_info, size=16, color=RGBColor(0xBB, 0xDE, 0xFB), space_before=Pt(12))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.8), Inches(4), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_ORANGE
    bar.line.fill.background()
    if num:
        slide_number(slide, num)
    return slide


def content_slide(title, bullets, num=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    top_bar(slide)
    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9))
    set_text(tb.text_frame, title, size=28, bold=True, color=DARK_BLUE)
    cb = add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    cb.text_frame.word_wrap = True
    first = True
    for item in bullets:
        if first:
            p = cb.text_frame.paragraphs[0]
            first = False
        else:
            p = cb.text_frame.add_paragraph()
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p.text = text
        p.font.size = Pt(opts.get("size", 18))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK_GRAY)
        p.font.italic = opts.get("italic", False)
        p.font.name = "Calibri"
        p.space_before = Pt(opts.get("space_before", 6))
        p.space_after = Pt(opts.get("space_after", 4))
        if opts.get("indent"):
            p.level = opts["indent"]
    if num:
        slide_number(slide, num)
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def exercise_slide(title, instructions, num=None, time_label="", note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, EXERCISE_BG)
    top_bar(slide, EXERCISE_ACCENT, 0.08)
    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    add_run(p, "EXERCISE  ", size=14, bold=True, color=EXERCISE_ACCENT)
    if time_label:
        add_run(p, time_label, size=14, color=MED_GRAY)
    add_paragraph(tb.text_frame, title, size=26, bold=True, color=DARK_BLUE, space_before=Pt(4))
    cb = add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
    cb.text_frame.word_wrap = True
    first = True
    for item in instructions:
        if first:
            p = cb.text_frame.paragraphs[0]
            first = False
        else:
            p = cb.text_frame.add_paragraph()
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p.text = text
        p.font.size = Pt(opts.get("size", 18))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK_GRAY)
        p.font.italic = opts.get("italic", False)
        p.font.name = "Calibri"
        p.space_before = Pt(opts.get("space_before", 6))
        p.space_after = Pt(opts.get("space_after", 4))
    if num:
        slide_number(slide, num)
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def debrief_slide(title, points, num=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DEBRIEF_BG)
    top_bar(slide, ACCENT_ORANGE, 0.08)
    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(1))
    p = tb.text_frame.paragraphs[0]
    add_run(p, "DEBRIEF  ", size=14, bold=True, color=ACCENT_ORANGE)
    add_paragraph(tb.text_frame, title, size=26, bold=True, color=DARK_BLUE, space_before=Pt(4))
    cb = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.3))
    cb.text_frame.word_wrap = True
    first = True
    for item in points:
        if first:
            p = cb.text_frame.paragraphs[0]
            first = False
        else:
            p = cb.text_frame.add_paragraph()
        if isinstance(item, tuple):
            text, opts = item
        else:
            text, opts = item, {}
        p.text = text
        p.font.size = Pt(opts.get("size", 18))
        p.font.bold = opts.get("bold", False)
        p.font.color.rgb = opts.get("color", DARK_GRAY)
        p.font.italic = opts.get("italic", False)
        p.font.name = "Calibri"
        p.space_before = Pt(opts.get("space_before", 6))
        p.space_after = Pt(opts.get("space_after", 4))
    if num:
        slide_number(slide, num)
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def table_slide(title, headers, rows, num=None, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    top_bar(slide)
    tb = add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8))
    set_text(tb.text_frame, title, size=28, bold=True, color=DARK_BLUE)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_w = Inches(11.5)
    tbl_h = Inches(min(n_rows * 0.7, 5.5))
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.8), Inches(1.3), tbl_w, tbl_h)
    tbl = tbl_shape.table
    col_w = tbl_w // n_cols
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = MID_BLUE
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
                p.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r_idx % 2 == 0 else LIGHT_GRAY
    if num:
        slide_number(slide, num)
    if note:
        slide.notes_slide.notes_text_frame.text = note
    return slide


def transition_slide(text, num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, TRANSITION_GREEN)
    tb = add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2.5))
    set_text(tb.text_frame, text, size=24, bold=False, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)
    tb.text_frame.paragraphs[0].font.italic = True
    if num:
        slide_number(slide, num)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# BLOCK 1 — OPENING (Slides 1-5, ~25 min)
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 1: Welcome ──
make_title_slide(
    "EVALAC I\nIntroduction to M&E and Theory of Change",
    "CLEAR-LAB  |  IDB  |  OVE\nSession 1 of 5  |  4.5 hours",
    num=1
)

# ── Slide 2: How this week works ──
content_slide("How This Week Works: 5 Sessions, One Story", [
    ("Each session builds on the previous one. All five use the same case study.", {"bold": True, "size": 20}),
    ("", {}),
    ("Session 1: How is change intended to happen?  (Theory of Change)", {"size": 17}),
    ("Session 2: How do we measure it?  (Indicators & Results Matrix)", {"size": 17}),
    ("Session 3: What qualitative evidence do we need?  (Qualitative Methods)", {"size": 17}),
    ("Session 4: What quantitative evidence do we need?  (Quantitative Methods)", {"size": 17}),
    ("Session 5: How strong is our contribution claim?  (Contribution Analysis)", {"size": 17}),
    ("", {}),
    ("We build the Theory of Change today and carry it forward all week.", {"italic": True, "color": MID_BLUE}),
    ("Session 2 picks up exactly where Session 1 ends.", {"italic": True, "color": MID_BLUE}),
], num=2, note="Emphasize: all 5 sessions use the same case study. We build the ToC today and carry it forward all week.")

# ── Slide 3: Learning objectives ──
content_slide("What We Will Learn Today: Session 1 Objectives", [
    ("Central question: How is change intended to happen -- and how do we reason about it?", {"bold": True, "size": 19, "color": MID_BLUE}),
    ("", {}),
    ("1. Frame complex, multi-actor interventions as systems of change with multiple interacting factors", {}),
    ("2. Explain why theory-based evaluation is essential, and position the Theory of Change as the starting point", {}),
    ("3. Describe how change is expected to happen: results, causal mechanisms, assumptions, external influences", {}),
    ("4. Distinguish monitoring from evaluation, and explain the role of the ToC in both", {}),
    ("5. Formulate evaluation questions that translate the ToC into evaluative inquiry", {}),
], num=3, note="Read objectives aloud; do not explain at length. Participants will understand them as the session unfolds. Return to this slide at the end as a check.")

# ── Slide 4: Diagnostic survey ──
content_slide("Diagnostic Survey: 3 Quick Questions", [
    ("Activate prior knowledge. Baseline the group.", {"bold": True, "size": 16, "color": MED_GRAY}),
    ("", {}),
    ("Q1: Have you ever built or revised a Theory of Change?", {"bold": True}),
    ("    Yes  /  No  /  Tried, not sure it worked", {"size": 16, "color": MED_GRAY}),
    ("", {}),
    ("Q2: What evaluation approach have you used most?", {"bold": True}),
    ("    RCT  /  Quasi-experimental  /  Process tracing  /  Contribution analysis  /  Other", {"size": 16, "color": MED_GRAY}),
    ("", {}),
    ("Q3: What is your biggest M&E challenge right now?", {"bold": True}),
    ("    1 sentence in chat or sticky note", {"size": 16, "color": MED_GRAY}),
], num=4, note="Display results live if using Miro/Mentimeter. Briefly name the distribution. Do not turn Q3 into a discussion -- collect and set aside. Use Q3 responses later as real-world anchors.")

# ── Slide 5: Meet the case ──
s5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s5, WHITE)
top_bar(s5)
tb = add_text_box(s5, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8))
set_text(tb.text_frame, "Meet the Case: Creciendo Juntos (Growing Together)", size=28, bold=True, color=DARK_BLUE)

# Case vignette box
box = add_shape_bg(s5, RGBColor(0xE8, 0xF0, 0xFE), Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.2))
vtb = add_text_box(s5, Inches(1.2), Inches(1.5), Inches(10.8), Inches(3.8))
vtb.text_frame.word_wrap = True
set_text(vtb.text_frame,
    "THE CASE: CRECIENDO JUNTOS",
    size=20, bold=True, color=DARK_BLUE)
add_paragraph(vtb.text_frame, "", size=8)
add_paragraph(vtb.text_frame,
    "A seven-year, government-run early childhood development program serving 180 municipalities. "
    "Target group: children aged 0-3 from the poorest 40% of families. "
    "Reached group: ~85,000 caregivers (mothers, fathers, grandmothers).",
    size=17, color=DARK_GRAY, space_before=Pt(8))
add_paragraph(vtb.text_frame,
    "Four components: (1) Weekly parenting group sessions (36/year)  (2) Monthly home visits  "
    "(3) Growth monitoring + developmental screening at health centers  (4) Early learning centers",
    size=17, color=DARK_GRAY, space_before=Pt(8))
add_paragraph(vtb.text_frame,
    "Budget: USD 285M over 5 years (~USD 1,140/child/year). Delivered by 750 community facilitators.",
    size=17, color=DARK_GRAY, space_before=Pt(8))
add_paragraph(vtb.text_frame,
    "Central question: Does participation by caregivers improve child development outcomes by age 3?",
    size=18, bold=True, color=ACCENT_ORANGE, space_before=Pt(14))

# Narrative reference
ntb = add_text_box(s5, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.2))
p_nt = set_text(ntb.text_frame,
    "Narrative Handout: S1_Narrative_Handout (pre-reading distributed). Full narrative (Parts 1 & 2) assigned for Session 2.",
    size=14, color=MED_GRAY)
p_nt.font.italic = True
p_cs = add_paragraph(ntb.text_frame,
    "Cheat Sheet Reference: 06_CaseStudy_CheatSheet_Session1 -- overview + program design",
    size=14, italic=True, color=MED_GRAY)
slide_number(s5, 5)
s5.notes_slide.notes_text_frame.text = (
    "Read the vignette aloud (2 min). Pause on one question: 'Looking at this program -- if it works, "
    "how exactly does that happen? What has to change?' Hold the question; do not answer yet. "
    "Distribute 1-page vignette handout. Do NOT ask participants to read in class."
)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCK 2 — INTERVENTIONS AND THEIR RESULTS (Slides 6-13, ~55 min)
# ══════════════════════════════════════════════════════════════════════════════

make_section_slide("BLOCK 2", "Interventions and Their Results",
    "~55 minutes  |  Slides 6-13  |  Exercise 1", num=6)

# ── Slide 7: What is a public intervention? ──
table_slide("What Is a Public Intervention?", [
    "INPUTS", "ACTIVITIES", "OUTPUTS"
], [
    ["What we invest", "What we do", "What we deliver"],
    ["750 facilitators\nCurriculum\nUSD 285M", "Parenting group sessions\nHome visits\nHealth center visits", "85,000 caregivers reached\n36 sessions/group/year"],
], num=7, note="Key distinction: Outputs tell us what was done. They do NOT tell us whether anything changed. The program can deliver 36 sessions and still not change anything. This is why we need a Theory of Change.")

# ── Slide 8: What are results? ──
content_slide("What Are Results? Outcomes and Impact", [
    ("OUTCOME: A change in knowledge, attitude, behavior, or condition of a person or system resulting from the intervention", {"bold": True, "size": 17}),
    ("", {}),
    ("IMPACT: Long-term, higher-order change attributable (in part) to the intervention -- often at population level", {"bold": True, "size": 17}),
    ("", {}),
    ("In Creciendo Juntos:", {"bold": True, "color": MID_BLUE, "size": 19}),
    ("  Outcome = caregiver changes parenting behavior", {"size": 17}),
    ("  Impact = child achieves age-appropriate development by age 3", {"size": 17}),
    ("", {}),
    ("The program delivers sessions (output). Sessions are supposed to change caregivers (outcome). Changed caregivers are supposed to develop children (impact). But HOW does each step happen?", {"italic": True, "color": ACCENT_ORANGE, "size": 17}),
], num=8)

# ── Slide 9: The results chain ──
content_slide("The Results Chain: A Causal Map of Change", [
    ("INPUTS  -->  ACTIVITIES  -->  OUTPUTS  -->  OUTCOMES  -->  IMPACT", {"bold": True, "size": 20, "color": MID_BLUE}),
    ("", {}),
    ("INPUTS: 750 facilitators, curriculum, USD 285M", {"size": 16}),
    ("ACTIVITIES: Parenting sessions, home visits, health center visits, learning centers", {"size": 16}),
    ("OUTPUTS: 85,000 caregivers reached, 36 sessions/group/year", {"size": 16}),
    ("OUTCOMES: Caregiver behavior change (interaction quality, nutrition practices)", {"size": 16}),
    ("IMPACT: 65% of children achieve age-appropriate development by 24 months", {"size": 16}),
    ("", {}),
    ("Which of these is easiest to measure? Which is hardest to attribute to the program?", {"italic": True, "color": ACCENT_ORANGE}),
], num=9, note="Draw the chain on the board as you speak. The attribution question sets up the contribution analysis thread for Session 5.")

# ── Slide 10: Direct and distal outcomes ──
content_slide("Direct and Distal Outcomes: Why the Distinction Matters", [
    ("DIRECT OUTCOME: Change in the group directly reached by the program", {"bold": True, "size": 17}),
    ("  Caregivers gain knowledge, change behavior", {"size": 16, "color": MED_GRAY}),
    ("", {}),
    ("DISTAL OUTCOME: Change in the group ultimately targeted", {"bold": True, "size": 17}),
    ("  Children show improved development", {"size": 16, "color": MED_GRAY}),
    ("", {}),
    ("THE INDIRECT PATHWAY", {"bold": True, "size": 20, "color": ACCENT_ORANGE}),
    ("Program  -->  Caregiver (direct)  -->  Child (distal)", {"bold": True, "size": 18, "color": MID_BLUE}),
    ("", {}),
    ("The program does NOT directly develop children's brains.", {"size": 17}),
    ("It develops caregivers' capacity to develop children's brains.", {"size": 17, "bold": True}),
    ("Caregiver behavioral change is the critical mechanism in this Theory of Change.", {"italic": True, "color": ACCENT_ORANGE}),
], num=10, note="Cheat Sheet S1, Ex.1: The program REACHES caregivers to CHANGE caregivers, which then AFFECTS children. A ToC fails if it jumps from 'sessions delivered' to 'child development improves' without specifying the behavioral mechanism.")

# ── Slide 11: Monitoring and evaluation brief ──
table_slide("Monitoring and Evaluation: Different Questions, Shared Theory", [
    "", "MONITORING", "EVALUATION"
], [
    ["Question", "Are we delivering as planned?", "Are outcomes occurring as expected?"],
    ["Focus", "Tracks outputs and early outcomes", "Tests whether the ToC holds -- and why"],
    ["Signal", "Early warning: is the theory unfolding?", "Strength of the contribution claim"],
    ["Timing", "Continuous -- during implementation", "Periodic -- at key milestones"],
], num=11, note="Both M&E use the same Theory of Change. That is why building a good ToC today is foundational to everything that follows this week.")

# ── Slide 12: Exercise 1 ──
exercise_slide("Exercise 1: Map the Results Chain for Creciendo Juntos", [
    ("15 min  |  Small groups (3-4 participants)  |  10 min groups + 5 min plenary", {"size": 14, "color": MED_GRAY, "bold": True}),
    ("", {}),
    ("Using the 1-page vignette distributed at Slide 5:", {"bold": True}),
    ("", {}),
    ("1. Complete the results chain template for Creciendo Juntos.", {}),
    ("   For each level, identify: (a) What is happening? (b) Who is involved? (c) Direct or distal change?", {"size": 16}),
    ("", {}),
    ("2. Identify ONE change that was probably NOT planned by the program.", {}),
    ("", {}),
    ("3. Be ready to share: What is the program trying to change, and in whom?", {}),
    ("", {}),
    ("Materials: 1-page case vignette + Results chain blank template (workbook)", {"size": 14, "color": MED_GRAY}),
    ("Cheat Sheet Ref: S1 Cheat Sheet, Exercise 1 -- Reached groups vs target groups", {"size": 14, "color": EXERCISE_ACCENT, "bold": True}),
], num=12, time_label="15 min",
note="ANSWER KEY (from Cheat Sheet S1):\n"
     "- Inputs: facilitators, curriculum, budget, health center staff, coordinators\n"
     "- Outputs: 85,000 caregivers reached, 36 sessions/group/year, 750 facilitators deployed\n"
     "- Outcomes (direct): Caregivers improve knowledge, skills, motivation -> change parenting behaviors\n"
     "- Outcomes (distal): Children 0-3 show improved language, motor, social-emotional development\n"
     "- Impact: Reduced developmental delay; long-term human capital gains\n"
     "- Unintended: Community facilitators gain professional credentials; peer networks reduce maternal isolation\n\n"
     "REACHED GROUPS (Cheat Sheet):\n"
     "Primary: Low-income caregivers (mothers primarily, some fathers, grandmothers) in 180 municipalities\n"
     "Secondary: Community facilitators, health center staff, municipal coordinators, educators\n"
     "Tertiary: Extended family, peer networks (indirect exposure)\n\n"
     "TARGET GROUPS:\n"
     "Primary: Children aged 0-5 (developmental outcomes)\n"
     "Secondary: Caregivers (knowledge, skills, well-being)")

# ── Slide 13: Debrief Ex 1 ──
debrief_slide("Debrief: The Indirect Pathway Is the Point", [
    ("Key insight: most groups correctly identify outputs but struggle with the mechanism between outputs and outcomes.", {}),
    ("", {}),
    ("The program REACHES caregivers to CHANGE caregivers, which then AFFECTS children.", {"bold": True, "color": DARK_BLUE}),
    ("", {}),
    ("Critical distinction:", {"bold": True}),
    ("  - The program delivers services to caregivers (direct)", {}),
    ("  - It expects children to benefit (distal)", {}),
    ("  - What is MISSING from the results chain: it tells us WHAT should happen, not WHY or HOW", {}),
    ("", {}),
    ("What you just built is a results chain. It is necessary -- but not sufficient.", {"italic": True, "color": ACCENT_ORANGE}),
    ("A Theory of Change adds three things: causal mechanisms, causal assumptions, and external influences.", {"italic": True, "color": ACCENT_ORANGE}),
], num=13, note="Bridge: 'We know WHAT should change. But WHY should it change -- and HOW exactly? That is what the Theory of Change needs to explain.'")

# ── Slide 14: Unintended outcomes ──
content_slide("Unintended Outcomes and Evaluation Breadth", [
    ("Programs regularly generate effects beyond the intended results chain.", {"size": 17}),
    ("", {}),
    ("A good Theory of Change flags plausible unintended pathways -- not to manage all of them, but to avoid being surprised.", {}),
    ("", {}),
    ("Examples from Creciendo Juntos:", {"bold": True, "color": MID_BLUE}),
    ("  - Facilitators develop professional credentials and career pathways", {"size": 16}),
    ("  - Peer networks among caregivers reduce maternal isolation and depression", {"size": 16}),
    ("  - Grandmothers exposed to new parenting norms change intergenerational practices", {"size": 16}),
    ("", {}),
    ("Unintended outcomes can be positive, negative, or mixed. Evaluation should scan for them.", {"italic": True}),
], num=14)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCK 3 — THEORY OF CHANGE (Slides 15-22, ~70 min)
# ══════════════════════════════════════════════════════════════════════════════

make_section_slide("BLOCK 3", "Theory of Change",
    "~70 minutes  |  Slides 15-22  |  Exercise 2", num=15)

# ── Slide 16: What is a Theory of Change? ──
content_slide("What Is a Theory of Change?", [
    ("DEFINITION", {"bold": True, "size": 22, "color": MID_BLUE}),
    ("A Theory of Change is a causal argument that explains how and why an intervention is expected to produce change.", {"size": 17}),
    ("", {}),
    ("It specifies:", {"bold": True}),
    ("  (1) The results chain", {}),
    ("  (2) The causal mechanisms", {}),
    ("  (3) The causal assumptions", {}),
    ("  (4) External influences", {}),
    ("", {}),
    ("What makes it different from a logic model:", {"bold": True, "color": ACCENT_ORANGE}),
    ("  - A logic model describes; a ToC argues", {"size": 16}),
    ("  - A logic model says 'if we do X, then Y'", {"size": 16}),
    ("  - A ToC says 'if we do X, then Y, BECAUSE OF Z, assuming A and B hold'", {"size": 16, "bold": True}),
    ("  - A ToC makes mechanism and assumptions explicit -- and therefore testable", {"size": 16}),
    ("", {}),
    ("Rogers (2008); Koleros & Mayne (2017)", {"size": 13, "color": MED_GRAY, "italic": True}),
], num=16)

# ── Slide 17: What is missing from the results chain? ──
content_slide("What Is Missing from the Results Chain?", [
    ("Results chain says:", {"bold": True, "color": MID_BLUE}),
    ("Sessions delivered  -->  Caregiver behavior change  -->  Child development", {"size": 17}),
    ("", {}),
    ("But does NOT explain:", {"bold": True, "color": ACCENT_ORANGE, "size": 20}),
    ("", {}),
    ("WHY would a session change caregiver behavior?", {"bold": True}),
    ("  --> Causal mechanism", {"color": MID_BLUE, "size": 16}),
    ("", {}),
    ("WHAT must be true for that change to happen?", {"bold": True}),
    ("  --> Causal assumptions", {"color": MID_BLUE, "size": 16}),
    ("", {}),
    ("WHAT else could produce the observed outcome?", {"bold": True}),
    ("  --> External influences", {"color": MID_BLUE, "size": 16}),
    ("", {}),
    ("The jump from 'sessions delivered' to 'child development improved' is not self-evident. It is a theory. And like any theory, it can be wrong.", {"italic": True, "color": ACCENT_ORANGE, "size": 17}),
], num=17)

# ── Slide 18: Causal mechanisms ──
content_slide("Causal Mechanisms: How Does Change Actually Happen?", [
    ("A causal mechanism is the process through which an output produces an outcome.", {"bold": True, "size": 17}),
    ("", {}),
    ("Key distinctions:", {"bold": True, "color": MID_BLUE}),
    ("  - Mechanisms are NOT activities", {}),
    ("  - Activities produce outputs; mechanisms explain how outputs produce outcomes", {}),
    ("  - In theory-based evaluation, we test mechanisms directly -- not just outcomes", {}),
    ("", {}),
    ("EXAMPLE: CRECIENDO JUNTOS", {"bold": True, "color": ACCENT_ORANGE, "size": 19}),
    ("  Activity: 36 parenting group sessions delivered", {"size": 16}),
    ("  Mechanism: Sessions increase caregiver CAPABILITY (knowledge), create OPPORTUNITY (peer networks, materials), and build MOTIVATION (new identity)", {"size": 16}),
    ("  Outcome: Caregivers change daily interaction behaviors", {"size": 16}),
    ("  Distal outcome: Children experience improved stimulation --> developmental gains", {"size": 16}),
    ("", {}),
    ("If the mechanism does not fire, sessions may be delivered but nothing changes.", {"italic": True, "color": ACCENT_ORANGE}),
], num=18)

# ── Slide 19: COM-B framework ──
table_slide("COM-B: Capability, Opportunity, Motivation --> Behavior", [
    "Component", "Meaning", "In Creciendo Juntos"
], [
    ["CAPABILITY (C)", "What a person knows and can do",
     "Caregiver learns milestones, learns responsive play"],
    ["OPPORTUNITY (O)", "What the environment enables or constrains",
     "Peer group forms, materials distributed, health visits accessible"],
    ["MOTIVATION (M)", "What drives action",
     "Child responds positively; new identity: 'I am my child's first teacher'"],
    ["BEHAVIOR (B)", "Observable action that produces the outcome",
     "Daily interaction increases 15 -> 45 min; harsh discipline replaced"],
], num=19, note="COM-B is an analytical lens, not a checklist. Its value is forcing the question: what exactly needs to shift in the person for new behavior to become possible? Cheat Sheet S1, Exercise 2 uses COM-B dimensions for caregiver behaviors.")

# ── Slide 20: COM-B applied ──
content_slide("COM-B in Creciendo Juntos: Tracing the Pathway", [
    ("CJ sessions --> [C] + [O] + [M] --> Behavior change --> Child development", {"bold": True, "size": 18, "color": MID_BLUE}),
    ("", {}),
    ("CAPABILITY:", {"bold": True, "color": ACCENT_GREEN}),
    ("  Knowledge of child development stages | Skill in responsive interaction | Self-efficacy", {"size": 16}),
    ("", {}),
    ("OPPORTUNITY:", {"bold": True, "color": ACCENT_GREEN}),
    ("  Peer support group | Learning materials | Health center access", {"size": 16}),
    ("", {}),
    ("MOTIVATION:", {"bold": True, "color": ACCENT_GREEN}),
    ("  Emotional reward from child's positive response | Peer norms | New parental identity", {"size": 16}),
    ("", {}),
    ("BEHAVIOR:", {"bold": True, "color": ACCENT_ORANGE}),
    ("  Responsive interaction UP | Positive discipline | Enriched home environment", {"size": 16}),
    ("", {}),
    ("We traced the mechanism for one actor -- the caregiver. In the exercise, you will go deeper.", {"italic": True, "color": MID_BLUE}),
], num=20, note="Cheat Sheet S1, Exercise 2: Defines observable behaviors for caregivers using COM-B. Also covers facilitators and health center staff as secondary actors.")

# ── Slide 21: Exercise 2 ──
exercise_slide("Exercise 2: Build the Causal Pathway for Creciendo Juntos", [
    ("20 min  |  Pairs or small groups  |  15 min groups + 5 min plenary", {"size": 14, "color": MED_GRAY, "bold": True}),
    ("", {}),
    ("Using the case and the COM-B lens, choose ONE reached group:", {"bold": True}),
    ("  (a) Caregivers    OR    (b) Community health workers / facilitators", {}),
    ("", {}),
    ("For your chosen group, trace the complete causal pathway:", {}),
    ("  (a) CAPABILITY: What must this person know or be able to do?", {}),
    ("  (b) OPPORTUNITY: What environmental conditions must be present?", {}),
    ("  (c) MOTIVATION: What must shift in their sense of purpose or habit?", {}),
    ("  (d) BEHAVIOR: What specific observable action should change?", {}),
    ("  (e) OUTCOME: What change does this behavior produce in the target group?", {}),
    ("", {}),
    ("Add your pathway to the results chain template from Exercise 1.", {"bold": True}),
    ("Cheat Sheet Ref: S1 Cheat Sheet, Exercise 2 -- Observable behaviors by actor (COM-B)", {"size": 14, "color": EXERCISE_ACCENT, "bold": True}),
], num=21, time_label="20 min",
note="ANSWER KEY -- Caregivers (from Cheat Sheet S1, Ex. 2):\n"
     "Capability: Knowledge of 6+ developmental milestones, responsive interaction skills, confidence\n"
     "Opportunity: Play materials, designated play space, contact with 3+ peers outside sessions, health services\n"
     "Motivation: Prioritizes daily interaction despite demands, uses positive discipline, shares learning\n"
     "Behavior: Daily interaction to 45 min, language-rich activities, replaces harsh discipline\n"
     "Outcome: Child receives stimulation -> developmental gains by 24 months\n\n"
     "ANSWER KEY -- Facilitators:\n"
     "Capability: Models responsive interaction, adapts curriculum to local context\n"
     "Opportunity/Motivation: Attends supervision, implements feedback, identifies/refers families with special needs\n\n"
     "ANSWER KEY -- Health Center Staff:\n"
     "Integrates developmental guidance into growth monitoring, administers ASQ-3, provides nutrition counseling")

# ── Slide 22: Debrief Ex 2 ──
debrief_slide("Debrief: Connecting Outputs --> Behaviors --> Outcomes", [
    ("The behavioral mechanism is what sessions are supposed to trigger.", {}),
    ("Without it, outputs do not produce outcomes.", {}),
    ("", {}),
    ("Key points:", {"bold": True, "color": DARK_BLUE}),
    ("  - Different groups trace different pathways -- both matter", {}),
    ("  - The pathway you traced is a hypothesis. Each link depends on something being true.", {}),
    ("  - Teaching note: Capability builds first (months 0-3), behavior emerges (months 3-6),", {"size": 16}),
    ("    consolidates (months 6-9), child outcomes appear (months 9-12+).", {"size": 16}),
    ("", {}),
    ("You have built the core of the Theory of Change.", {"bold": True, "color": ACCENT_ORANGE}),
    ("The next question: what are we assuming? What must be true for each link to hold?", {"italic": True, "color": ACCENT_ORANGE}),
], num=22, note="Bridge: 'We have been assuming a lot. What must be true for this pathway to hold? That is what causal assumptions tell us.'")

# ── Slide 23: Full ToC reveal ──
content_slide("The Creciendo Juntos ToC: The Full Causal Picture", [
    ("This is the complete Theory of Change incorporating all elements from Exercises 1 & 2.", {"size": 16, "color": MED_GRAY}),
    ("", {}),
    ("Layer 1 (BLACK): Results chain -- from Exercise 1", {"bold": True, "size": 17}),
    ("  Inputs --> Activities --> Outputs --> Outcomes --> Impact", {"size": 16}),
    ("", {}),
    ("Layer 2 (BLUE): COM-B behavioral mechanism -- from Exercise 2", {"bold": True, "size": 17, "color": MID_BLUE}),
    ("  Capability + Opportunity + Motivation --> Behavior change", {"size": 16, "color": MID_BLUE}),
    ("", {}),
    ("Layer 3 (ORANGE): Assumption markers -- to be filled in Block 4", {"bold": True, "size": 17, "color": ACCENT_ORANGE}),
    ("  Conditions that must hold for each link", {"size": 16, "color": ACCENT_ORANGE}),
    ("", {}),
    ("Layer 4 (GREY): External influence arrows -- to be filled in Block 4", {"bold": True, "size": 17, "color": MED_GRAY}),
    ("  Factors outside the program's control", {"size": 16, "color": MED_GRAY}),
    ("", {}),
    ("What surprises you? What would you add or question?", {"italic": True, "color": ACCENT_ORANGE}),
], num=23, note="Reveal full annotated ToC from Cheat Sheet S1 (File 06). Ask: 'What surprises you? What would you add or question?' Take 2-3 responses. Keep brief -- Block 4 will go deep on assumptions.")

# ══════════════════════════════════════════════════════════════════════════════
# BLOCK 4 — ASSUMPTIONS AND EXTERNAL INFLUENCES (Slides 24-31, ~65 min)
# ══════════════════════════════════════════════════════════════════════════════

make_section_slide("BLOCK 4", "Assumptions and External Influences",
    "~65 minutes  |  Slides 24-31  |  Exercises 3 & 4", num=24)

# ── Slide 25: What are causal assumptions? ──
content_slide("What Are Causal Assumptions?", [
    ("DEFINITION", {"bold": True, "size": 20, "color": MID_BLUE}),
    ("A causal assumption is a condition that must hold for one level of the results chain to produce the next.", {}),
    ("If an assumption is violated, the expected causal link breaks -- even if the program is implemented correctly.", {}),
    ("", {}),
    ("Example from Creciendo Juntos:", {"bold": True, "color": ACCENT_ORANGE}),
    ("  Link: Parenting sessions delivered --> Caregivers change behavior", {"size": 16}),
    ("  Assumption: Attendance is sufficient (caregivers attend >= 70% of sessions)", {"size": 16}),
    ("", {}),
    ("If violated:", {"bold": True, "color": RGBColor(0xC6, 0x28, 0x28)}),
    ("  Sessions delivered but caregivers attend sporadically --> capability does not build", {"size": 16}),
    ("  --> behavior does not change --> outcomes not reached", {"size": 16}),
    ("  Even though outputs were 100% delivered.", {"size": 16, "bold": True}),
    ("", {}),
    ("This is why monitoring attendance is not optional. If this assumption breaks, everything downstream breaks.", {"italic": True, "color": ACCENT_ORANGE}),
], num=25)

# ── Slide 26: Types of assumptions ──
table_slide("Types of Assumptions: Classifying by Results Chain Level", [
    "Type", "Question It Answers", "Creciendo Juntos Example"
], [
    ["REACH", "Can the program reach the intended population?",
     "Facilitators can be recruited and retained in target communities"],
    ["CAPACITY", "Does the intervention build the required COM-B shift?",
     "Sessions delivered with sufficient quality to build caregiver capability"],
    ["BEHAVIOR", "Does changed COM-B produce behavior change?",
     "Knowledge and motivation translate to daily behavior (not just during sessions)"],
    ["OUTCOME", "Does behavior change produce the expected outcome?",
     "Changed caregiver behavior sustained 18+ months produces developmental gains"],
], num=26, note="Cheat Sheet S1, Exercise 4: Uses format 'IF [this happens], THEN [that follows], ASSUMING [critical condition].' Assumptions classified by results chain link.")

# ── Slide 27: Critical assumptions ──
content_slide("Critical Assumptions: Which Links Are Most Fragile?", [
    ("A critical assumption is one that:", {"bold": True, "size": 19}),
    ("  1. Is ESSENTIAL -- if it breaks, outcomes cannot be reached", {}),
    ("  2. Is UNCERTAIN -- we have limited evidence it holds in this context", {}),
    ("  3. Is TESTABLE -- we can design monitoring or evaluation to check it", {}),
    ("", {}),
    ("TWO CANDIDATE CRITICAL ASSUMPTIONS:", {"bold": True, "color": ACCENT_ORANGE, "size": 19}),
    ("", {}),
    ("B1 (BEHAVIOR):", {"bold": True, "color": MID_BLUE}),
    ("Caregivers who increase knowledge in sessions also change their daily parenting behaviors at home, not just during observed sessions.", {"size": 16}),
    ("CRITICAL because behavioral generalization is uncertain and rarely measured.", {"size": 15, "italic": True, "color": MED_GRAY}),
    ("", {}),
    ("O1 (OUTCOME):", {"bold": True, "color": MID_BLUE}),
    ("Changed parenting behavior, sustained for 18 months, is sufficient to produce measurable developmental gains in children by age 3.", {"size": 16}),
    ("CRITICAL because the dose-response relationship is unclear for multi-risk households.", {"size": 15, "italic": True, "color": MED_GRAY}),
], num=27)

# ── Slide 28: Exercise 3 ──
exercise_slide("Exercise 3: Identify and Classify Assumptions in the ToC", [
    ("20 min  |  Small groups (3-4)  |  15 min groups + 5 min plenary", {"size": 14, "color": MED_GRAY, "bold": True}),
    ("", {}),
    ("Using the Creciendo Juntos ToC from Slide 23:", {"bold": True}),
    ("", {}),
    ("1. Identify at least ONE assumption per level: reach, capacity, behavior, outcome.", {}),
    ("", {}),
    ("2. Write each assumption as a complete sentence:", {}),
    ("   'IF [condition], THEN the link holds.'", {"italic": True, "size": 16}),
    ("", {}),
    ("3. Mark which TWO assumptions you consider most critical.", {}),
    ("", {}),
    ("4. Be ready to explain: Why these two? What evidence would test them?", {}),
    ("", {}),
    ("Add assumptions to the ToC template -- they annotate each causal arrow.", {"bold": True}),
    ("Cheat Sheet Ref: S1 Cheat Sheet, Exercise 4 -- Assumptions by causal link", {"size": 14, "color": EXERCISE_ACCENT, "bold": True}),
], num=28, time_label="20 min",
note="ANSWER KEY (from Cheat Sheet S1, Ex. 4):\n"
     "Link 1 (Inputs -> Outputs):\n"
     "  A1: Facilitators can be recruited, trained, and retained\n"
     "  A2: Municipalities provide space and logistical support\n"
     "  A3: Health centers have capacity to add services\n\n"
     "Link 2 (Outputs -> Behavior Change) -- MOST CRITICAL:\n"
     "  A4: Caregivers attend regularly (>=70%). Risk: poverty, transport, family opposition\n"
     "  A5: Knowledge/skills translate to daily practice. Risk: knowing != doing\n"
     "  A6: Changed behavior sustains after program ends. Risk: practices fade\n\n"
     "Link 3 (Behavior Change -> Outcomes):\n"
     "  A7: Dosage and quality of changed interaction is sufficient\n"
     "  A8: No overwhelming adversities swamp parenting improvements\n\n"
     "Bridge: 'The assumptions you identified are exactly the questions a contribution analysis must test. Session 5 will ask: which held? How strong is our claim?'")

# ── Slide 29: Debrief Ex 3 ──
debrief_slide("Debrief: Which Assumptions Are Most Fragile?", [
    ("Most groups identify Behavior assumptions as most critical -- this is where most ECD programs fail.", {}),
    ("", {}),
    ("An assumption is only useful if it is testable.", {"bold": True}),
    ("", {}),
    ("Critical assumptions become:", {"bold": True, "color": DARK_BLUE}),
    ("  - Monitoring indicators (track whether conditions hold)", {}),
    ("  - Evaluation questions (test whether the ToC mechanism functioned)", {}),
    ("", {}),
    ("Example:", {"bold": True, "color": ACCENT_ORANGE}),
    ("If the critical assumption is 'caregivers change daily behavior,' the evaluation question is:", {"size": 16}),
    ("Did caregivers actually change their daily parenting behavior -- and if so, how much, for whom, and under what conditions?", {"size": 16, "italic": True}),
    ("", {}),
    ("That is what Session 5 will help you answer.", {"color": MID_BLUE, "italic": True}),
], num=29)

# ── Slide 30: External influences ──
content_slide("External Influences: What Is Outside the Program's Control?", [
    ("Two types of external influence:", {"bold": True, "size": 20, "color": MID_BLUE}),
    ("", {}),
    ("CONTEXTUAL FACTORS", {"bold": True, "color": ACCENT_GREEN}),
    ("Structural conditions that shape whether the mechanism can function:", {}),
    ("  Economy, political stability, health system quality, food security", {"size": 16}),
    ("", {}),
    ("RIVAL EXPLANATIONS", {"bold": True, "color": ACCENT_ORANGE}),
    ("Alternative causes of the observed outcome without the program:", {}),
    ("  Natural child development, another program, a social trend", {"size": 16}),
    ("", {}),
    ("Why both matter for evaluation:", {"bold": True}),
    ("  - Contextual factors can violate assumptions (e.g., severe food insecurity means even improved parenting doesn't translate to gains)", {"size": 16}),
    ("  - Rival explanations: if children improve even in control communities, the program's contribution is less clear", {"size": 16}),
    ("  - Evaluators must rule out rival explanations or account for them in the contribution claim", {"size": 16}),
], num=30)

# ── Slide 31: Exercise 4 ──
exercise_slide("Exercise 4: External Influences Scan (Facilitated Discussion)", [
    ("10 min  |  Full-group facilitated discussion  |  No separate template needed", {"size": 14, "color": MED_GRAY, "bold": True}),
    ("", {}),
    ("Q1: What are the 2-3 most important external factors that could affect whether Creciendo Juntos achieves its outcomes?", {"bold": True}),
    ("  (Facilitator records on board)", {"size": 14, "color": MED_GRAY}),
    ("", {}),
    ("Q2: Which of these could be a rival explanation for observed improvements?", {"bold": True}),
    ("  Which could make outcomes improve even if the program had no effect?", {}),
    ("", {}),
    ("Q3: Which critical assumption from Exercise 3 is most vulnerable to these external factors?", {"bold": True}),
    ("", {}),
    ("Expected factors: economic recession, political instability, other social programs (CCTs),", {"size": 15, "italic": True, "color": MED_GRAY}),
    ("natural disasters / pandemics, media campaigns on parenting", {"size": 15, "italic": True, "color": MED_GRAY}),
    ("Cheat Sheet Ref: S1 Cheat Sheet, Exercise 5 -- ToC to Contribution Analysis preview", {"size": 14, "color": EXERCISE_ACCENT, "bold": True}),
], num=31, time_label="10 min",
note="Bridge: 'These factors go on the edge of the ToC diagram -- they are not inside the program's theory, but they affect whether it holds. In evaluation design, we call these threats to internal validity. We return to them in Sessions 3-5.'")

# ── Slide 32: Complete ToC - four layers ──
content_slide("The Complete Theory of Change: Four Layers", [
    ("You have now built the full Creciendo Juntos Theory of Change.", {"bold": True, "size": 19, "color": MID_BLUE}),
    ("", {}),
    ("Layer 1 (BLACK): Results chain -- from Exercise 1", {"bold": True}),
    ("  Inputs --> Activities --> Outputs --> Outcomes --> Impact", {"size": 16}),
    ("", {}),
    ("Layer 2 (BLUE): Behavioral mechanism (COM-B) -- from Exercise 2", {"bold": True, "color": MID_BLUE}),
    ("  Capability + Opportunity + Motivation --> Behavior change", {"size": 16}),
    ("", {}),
    ("Layer 3 (ORANGE): Causal assumptions, classified by level -- from Exercise 3", {"bold": True, "color": ACCENT_ORANGE}),
    ("  Reach | Capacity | Behavior | Outcome assumptions", {"size": 16}),
    ("", {}),
    ("Layer 4 (GREY): External influence arrows -- from Exercise 4", {"bold": True, "color": MED_GRAY}),
    ("  Contextual factors and rival explanations", {"size": 16}),
    ("", {}),
    ("This is a Theory of Change. It is not a diagram -- it is an argument.", {"bold": True, "italic": True, "color": ACCENT_ORANGE, "size": 19}),
    ("You can now use it to design evaluation questions, select indicators, and assess your contribution claim.", {"italic": True, "size": 17}),
], num=32)

# ══════════════════════════════════════════════════════════════════════════════
# BLOCK 5 — M&E + BRIDGE TO SESSION 2 (Slides 33-41, ~30 min)
# ══════════════════════════════════════════════════════════════════════════════

make_section_slide("BLOCK 5", "M&E and Bridge to Session 2",
    "~30 minutes  |  Slides 33-41  |  Exercise 5 (optional)", num=33)

# ── Slide 34: ToC and evaluation questions ──
table_slide("ToC and Evaluation Questions: Making the Theory Testable", [
    "Critical Assumption (from Ex. 3)", "Evaluation Question It Generates"
], [
    ["Behavior: Caregivers change daily practices,\nnot just session behavior",
     "Did caregivers change their daily interaction\nbehavior at home? For whom? To what extent?"],
    ["Outcome: Changed behavior sustained 18 months\nproduces developmental gains",
     "Is there a relationship between intensity of\nbehavior change and child outcomes?"],
    ["Reach: Facilitators retained in communities\nfor program duration",
     "What is facilitator turnover rate and how does\nit affect session continuity and caregiver trust?"],
], num=34, note="Key teaching point: 'You did not just build a diagram. You built a set of evaluation questions. Evaluation design starts here -- from the assumptions.'")

# ── Slide 35: Monitoring ──
content_slide("Monitoring: Watching the ToC Unfold in Real Time", [
    ("Monitoring observes: reach, outputs, and early behavioral signals.", {}),
    ("", {}),
    ("Monitoring does NOT test the ToC -- it signals whether conditions for the ToC to function are present.", {"bold": True}),
    ("", {}),
    ("Example:", {"bold": True, "color": ACCENT_ORANGE}),
    ("  If monitoring shows low attendance --> reach assumption may be breaking", {"size": 16}),
    ("  --> evaluator can investigate before outcomes are affected", {"size": 16}),
    ("", {}),
    ("Connection to Session 2:", {"bold": True, "color": MID_BLUE}),
    ("What you monitor is determined by your ToC.", {}),
    ("In Session 2, we take this results chain and ask:", {}),
    ("  What indicator would tell us whether each level is happening?", {"italic": True}),
    ("  That is the Results Matrix.", {"italic": True, "bold": True}),
], num=35)

# ── Slide 36: Voices from the field (Narrative integration) ──
s36 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s36, QUOTE_BG)
top_bar(s36, RGBColor(0x7B, 0x1F, 0xA2), 0.08)
tb36 = add_text_box(s36, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8))
set_text(tb36.text_frame, "Voices from the Field: Why the ToC Matters", size=28, bold=True, color=DARK_BLUE)
cb36 = add_text_box(s36, Inches(1.2), Inches(1.4), Inches(10.8), Inches(5.5))
cb36.text_frame.word_wrap = True
p_q1 = set_text(cb36.text_frame,
    '"I love my daughter, but I honestly do not know what to do with her all day." -- Maria, 22',
    size=17, color=DARK_GRAY)
p_q1.font.italic = True
add_paragraph(cb36.text_frame,
    '"I am alone all day with three children under five." -- Carmen, 27',
    size=17, italic=True, color=DARK_GRAY, space_before=Pt(14))
add_paragraph(cb36.text_frame,
    '"These young mothers do not know how to discipline. A smack teaches respect." -- Dona Rosa, 58',
    size=17, italic=True, color=DARK_GRAY, space_before=Pt(14))
add_paragraph(cb36.text_frame,
    '"My job is to provide money. I would be the only man there." -- Diego, 28',
    size=17, italic=True, color=DARK_GRAY, space_before=Pt(14))
add_paragraph(cb36.text_frame,
    '"The health center makes me feel stupid. The nurse says read to my baby. I do not have books." -- Valeria, 15',
    size=17, italic=True, color=DARK_GRAY, space_before=Pt(14))
add_paragraph(cb36.text_frame, "", size=8)
add_paragraph(cb36.text_frame,
    "Mendoza's synthesis: the challenge was not merely lack of services but a complex intersection of "
    "limited knowledge, conflicting beliefs, social isolation, structural barriers, and rigid gender norms.",
    size=16, bold=True, color=RGBColor(0x4A, 0x14, 0x8C), space_before=Pt(18))
add_paragraph(cb36.text_frame,
    "Source: S1_Narrative_Handout -- Family Listening Sessions",
    size=13, italic=True, color=MED_GRAY, space_before=Pt(14))
slide_number(s36, 36)
s36.notes_slide.notes_text_frame.text = (
    "These quotes come from the family listening sessions in the narrative. They illustrate WHY the ToC must specify "
    "behavioral mechanisms: each quote reveals a different barrier (knowledge, isolation, norms, gender, literacy). "
    "The ToC must account for ALL of these through COM-B."
)

# ── Slide 37: Exercise 5 (optional) ──
exercise_slide("Exercise 5: What Would You Monitor First? (Optional)", [
    ("8 min  |  Pairs  |  5 min pairs + 3 min share-out  |  Can be homework", {"size": 14, "color": MED_GRAY, "bold": True}),
    ("", {}),
    ("Looking at the completed Creciendo Juntos ToC:", {"bold": True}),
    ("", {}),
    ("1. Choose the THREE variables you would monitor most closely in the first 6 months.", {}),
    ("", {}),
    ("2. For each: (a) What would you measure?  (b) Who would you observe or survey?  (c) When?", {}),
    ("", {}),
    ("3. Explain: Why these three? What would each tell you about whether the ToC is unfolding?", {}),
    ("", {}),
    ("Note: These are proto-indicators. In Session 2, we formalize these into a Results Matrix.", {"italic": True, "color": MID_BLUE}),
    ("Cheat Sheet Ref: S1 Cheat Sheet, Exercise 5 -- Connecting ToC to Contribution Analysis", {"size": 14, "color": EXERCISE_ACCENT, "bold": True}),
], num=37, time_label="8 min (optional)",
note="Expected answers:\n"
     "- Facilitator turnover rate (reach assumption) -- monthly from HR records\n"
     "- Session attendance rate among caregivers (capacity assumption) -- per group per session\n"
     "- Observed caregiver interaction during home visit (behavior assumption) -- structured observation at Month 3 and 6")

# ── Slide 38: Evaluation testing ──
content_slide("Evaluation: Testing the ToC at Outcome Level", [
    ("Evaluation asks:", {"bold": True, "color": MID_BLUE}),
    ("  Did expected outcomes occur?", {}),
    ("  Did the theory's mechanisms function as hypothesized?", {}),
    ("", {}),
    ("Contribution analysis (Session 5):", {"bold": True, "color": ACCENT_ORANGE}),
    ("'Given what we observe, how strong is our claim that the program contributed to these outcomes?'", {"italic": True}),
    ("", {}),
    ("Evidence we can test (from Cheat Sheet S1):", {"bold": True}),
    ("  - Knowledge surveys (did capability increase?)", {"size": 16}),
    ("  - Peer network surveys (did opportunity expand?)", {"size": 16}),
    ("  - Time-use data (did behavior actually change at home?)", {"size": 16}),
    ("  - Observation rubrics (quality of caregiver-child interaction)", {"size": 16}),
    ("", {}),
    ("The ToC built today is the foundation of the contribution claim.", {"bold": True, "color": MID_BLUE}),
    ("Every assumption identified is a piece of evidence to test.", {}),
], num=38, note="Thread to Sessions 3-5: Sessions 3 and 4 will ask what qualitative and quantitative evidence lets you test each mechanism. Session 5 brings all evidence into a contribution claim.")

# ── Slide 39: From ToC to Results Matrix ──
s39 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s39, TRANSITION_GREEN)
top_bar(s39, ACCENT_GREEN, 0.08)
tb39 = add_text_box(s39, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8))
set_text(tb39.text_frame, "From ToC to Results Matrix: What Comes Next", size=28, bold=True, color=DARK_BLUE)

cb39 = add_text_box(s39, Inches(0.8), Inches(1.3), Inches(5.3), Inches(5.5))
cb39.text_frame.word_wrap = True
set_text(cb39.text_frame, "TODAY YOU BUILT:", size=20, bold=True, color=ACCENT_GREEN)
add_paragraph(cb39.text_frame, "A results chain (Exercise 1)", size=17, space_before=Pt(10))
add_paragraph(cb39.text_frame, "A causal mechanism layer -- COM-B (Exercise 2)", size=17, space_before=Pt(6))
add_paragraph(cb39.text_frame, "A classified assumption list (Exercise 3)", size=17, space_before=Pt(6))
add_paragraph(cb39.text_frame, "An external influence map (Exercise 4)", size=17, space_before=Pt(6))

cb39r = add_text_box(s39, Inches(6.5), Inches(1.3), Inches(6.3), Inches(5.5))
cb39r.text_frame.word_wrap = True
set_text(cb39r.text_frame, "SESSION 2 WILL ASK:", size=20, bold=True, color=MID_BLUE)
add_paragraph(cb39r.text_frame, "What indicator tells us whether each level is happening?", size=17, space_before=Pt(10))
add_paragraph(cb39r.text_frame, "What source? What frequency? What baseline?", size=17, space_before=Pt(6))
add_paragraph(cb39r.text_frame, "", size=8)
add_paragraph(cb39r.text_frame, "PREVIEW -- One Row of the Results Matrix:", size=16, bold=True, color=ACCENT_ORANGE, space_before=Pt(12))
add_paragraph(cb39r.text_frame, "Level: Caregiver behavior change", size=15, space_before=Pt(8))
add_paragraph(cb39r.text_frame, "Indicator: % caregivers with >= 3 responsive interactions/hour", size=15, space_before=Pt(4))
add_paragraph(cb39r.text_frame, "Baseline: 15% at enrollment  |  Target: 65% at Month 9", size=15, space_before=Pt(4))
add_paragraph(cb39r.text_frame, "Source: Structured observation (Months 3, 6, 9)", size=15, space_before=Pt(4))
slide_number(s39, 39)

# ══════════════════════════════════════════════════════════════════════════════
# CLOSURE (Slides 40-44, ~15 min)
# ══════════════════════════════════════════════════════════════════════════════

make_section_slide("CLOSURE", "Takeaways and Next Steps",
    "~15 minutes  |  Slides 40-44", num=40)

# ── Slide 41: Key takeaways ──
content_slide("Key Takeaways: Returning to the Learning Objectives", [
    ("1. Frame complex interventions", {"bold": True, "color": MID_BLUE}),
    ("   Creciendo Juntos reaches caregivers to change children; the indirect pathway is the mechanism", {"size": 16}),
    ("", {}),
    ("2. Explain theory-based evaluation", {"bold": True, "color": MID_BLUE}),
    ("   The results chain alone is insufficient; mechanisms, assumptions, and external influences are essential", {"size": 16}),
    ("", {}),
    ("3. Describe a Theory of Change", {"bold": True, "color": MID_BLUE}),
    ("   Four layers: results chain, mechanisms, assumptions, external influences", {"size": 16}),
    ("", {}),
    ("4. Distinguish M&E", {"bold": True, "color": MID_BLUE}),
    ("   Monitoring watches the ToC unfold; evaluation tests whether it held", {"size": 16}),
    ("", {}),
    ("5. Formulate evaluation questions", {"bold": True, "color": MID_BLUE}),
    ("   Each critical assumption generates a specific evaluative question", {"size": 16}),
    ("", {}),
    ("Closing reflection: Think of a program you work with. What is the most critical assumption -- the one that, if violated, means the program could deliver everything planned and still fail?", {"italic": True, "color": ACCENT_ORANGE, "size": 16}),
], num=41)

# ── Slide 42: Pre-reading for Session 2 ──
s42 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s42, DARK_BLUE)
bottom_accent(s42)
tb42 = add_text_box(s42, Inches(1), Inches(1.5), Inches(11), Inches(5))
set_text(tb42.text_frame, "BEFORE SESSION 2 -- READ:", size=24, bold=True, color=WHITE)
add_paragraph(tb42.text_frame, "", size=8)
add_paragraph(tb42.text_frame,
    "S2_Narrative_Handout: Program Design and Implementation Structure",
    size=20, color=WHITE, bold=True, space_before=Pt(14))
add_paragraph(tb42.text_frame,
    "Focus on: How are the four components coordinated? What roles do facilitators, health staff, and coordinators play?",
    size=17, color=RGBColor(0xBB, 0xDE, 0xFB), space_before=Pt(12))
add_paragraph(tb42.text_frame, "", size=8)
add_paragraph(tb42.text_frame,
    "In Session 2, we will use this program design to populate a full Results Matrix -- one indicator set per results chain level.",
    size=17, color=RGBColor(0xBB, 0xDE, 0xFB), space_before=Pt(12))
add_paragraph(tb42.text_frame,
    "Bring your completed ToC templates from today's exercises.",
    size=18, bold=True, color=ACCENT_ORANGE, space_before=Pt(20))
add_paragraph(tb42.text_frame,
    "Estimated reading time: 15-20 minutes.",
    size=15, color=RGBColor(0x88, 0xAA, 0xCC), space_before=Pt(14))
slide_number(s42, 42)

# ── Slide 43: Closing survey ──
content_slide("Closing Survey", [
    ("Standard course satisfaction items (per OVE format)", {"size": 16, "color": MED_GRAY}),
    ("", {}),
    ("Formative assessment:", {"bold": True, "size": 20, "color": MID_BLUE}),
    ("", {}),
    ("Q1: In one sentence: what is the causal mechanism in the Creciendo Juntos Theory of Change?", {"bold": True}),
    ("", {}),
    ("Q2: What is the most critical assumption in a Theory of Change you currently work with?", {"bold": True}),
    ("", {}),
    ("", {}),
    ("These items feed directly into the EVALAC pre/post evaluation.", {"size": 14, "italic": True, "color": MED_GRAY}),
], num=43)

# ══════════════════════════════════════════════════════════════════════════════
# REFERENCE SLIDES (44-50, not presented in session)
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 44: Glossary ──
content_slide("Reference: Glossary of Key Concepts", [
    ("Theory of Change: Causal argument explaining how/why an intervention produces change", {"size": 15}),
    ("Results Chain: Linear sequence from inputs to impact", {"size": 15}),
    ("Causal Mechanism: Process through which an output produces an outcome", {"size": 15}),
    ("COM-B: Capability + Opportunity + Motivation --> Behavior (Michie et al.)", {"size": 15}),
    ("Causal Assumption: Condition that must hold for one results level to produce the next", {"size": 15}),
    ("Critical Assumption: Essential, uncertain, and testable assumption", {"size": 15}),
    ("Direct Outcome: Change in the group directly reached by the program", {"size": 15}),
    ("Distal Outcome: Change in the group ultimately targeted", {"size": 15}),
    ("External Influence: Factor outside program control affecting outcomes", {"size": 15}),
    ("Rival Explanation: Alternative cause of observed outcomes", {"size": 15}),
    ("Monitoring: Continuous tracking of outputs and early outcomes", {"size": 15}),
    ("Evaluation: Periodic assessment of whether the ToC held and why", {"size": 15}),
    ("Contribution Analysis: Method for assessing strength of program's contribution claim", {"size": 15}),
    ("Results Matrix: Framework linking each results level to indicators, sources, and targets", {"size": 15}),
], num=44, note="Reference slide -- not presented in session.")

# ── Slide 45: Full annotated ToC ──
content_slide("Reference: Creciendo Juntos -- Full Annotated Theory of Change", [
    ("INPUTS: 750 facilitators, curriculum, USD 285M, health center staff, coordinators", {"size": 14}),
    ("  --> ACTIVITIES: Parenting sessions (36/yr), home visits, growth monitoring, learning centers", {"size": 14}),
    ("    --> OUTPUTS: 85,000 caregivers, 36 sessions/group/yr, 750 facilitators deployed", {"size": 14}),
    ("", {}),
    ("COM-B MECHANISM (caregiver):", {"bold": True, "color": MID_BLUE, "size": 14}),
    ("  C: Knowledge of milestones + responsive interaction skills + self-efficacy", {"size": 13}),
    ("  O: Peer support + materials + health services access", {"size": 13}),
    ("  M: Child's positive response + peer norms + new parental identity", {"size": 13}),
    ("  B: Responsive interaction UP, positive discipline, enriched home environment", {"size": 13}),
    ("", {}),
    ("  --> OUTCOMES (direct): Caregiver behavior change sustained", {"size": 14}),
    ("    --> OUTCOMES (distal): Child development (language, motor, social-emotional)", {"size": 14}),
    ("      --> IMPACT: 65% age-appropriate development by 24 months (vs 35% baseline)", {"size": 14}),
    ("", {}),
    ("ASSUMPTIONS: Reach (recruitment/retention) | Capacity (session quality) | Behavior (knowledge-->action) | Outcome (dose-response)", {"size": 13, "color": ACCENT_ORANGE}),
    ("EXTERNAL: Economic conditions | Political stability | Other programs (CCTs) | Natural development", {"size": 13, "color": MED_GRAY}),
], num=45, note="Reference slide -- 4-layer annotated ToC from File 06 cheat sheet.")

# ── Slide 46: Results Matrix preview ──
table_slide("Reference: Results Matrix Preview -- One Completed Row", [
    "Element", "Content"
], [
    ["Results chain level", "Caregiver behavior change"],
    ["Indicator", "% of caregivers observed engaging in >= 3 responsive\ninteraction activities per hour during structured home visit"],
    ["Type", "Behavioral (mechanism) indicator"],
    ["Baseline", "15% of caregivers at enrollment"],
    ["Target", "65% at Month 9"],
    ["Source", "Structured observation at home visits (Months 3, 6, 9)"],
], num=46, note="Reference slide -- bridge to Session 2.")

# ── Slide 47: COM-B lens card ──
table_slide("Reference: COM-B Lens Card", [
    "Dimension", "Definition", "Guiding Question"
], [
    ["CAPABILITY", "What a person knows and can do\n(physical + psychological)",
     "What must this person know or be\nable to do that they currently cannot?"],
    ["OPPORTUNITY", "What the environment enables\nor constrains (physical + social)",
     "What conditions must be present\nfor the behavior to occur?"],
    ["MOTIVATION", "What drives action\n(reflective + automatic)",
     "What must shift in their sense of\npurpose, identity, or habit?"],
    ["BEHAVIOR", "Observable action that\nproduces the outcome",
     "What specific, measurable action\nshould change?"],
], num=47, note="Reference card for exercises. Based on Michie et al. (2011) COM-B model.")

# ── Slide 48: Assumption classification card ──
table_slide("Reference: Assumption Classification Card", [
    "Assumption Type", "Results Chain Link", "Format", "CJ Example"
], [
    ["REACH", "Can program reach\nintended population?",
     "IF [reach condition],\nTHEN [population accessed]",
     "Facilitators recruited and\nretained in 180 municipalities"],
    ["CAPACITY", "Does intervention build\nrequired COM-B shift?",
     "IF [quality condition],\nTHEN [COM-B shift occurs]",
     "Sessions delivered with quality\nsufficient to build capability"],
    ["BEHAVIOR", "Does changed COM-B\nproduce behavior?",
     "IF [COM-B present],\nTHEN [behavior changes]",
     "Knowledge + motivation translate\nto daily practice at home"],
    ["OUTCOME", "Does behavior produce\nexpected outcome?",
     "IF [behavior sustained],\nTHEN [outcome observed]",
     "18+ months of changed behavior\nproduces developmental gains"],
], num=48, note="Reference card for Exercise 3.")

# ── Slide 49: Key readings ──
content_slide("Reference: Key Readings", [
    ("Foundational:", {"bold": True, "color": MID_BLUE}),
    ("  Rogers, P. (2008). Using Programme Theory to Evaluate Complicated and Complex Aspects of Interventions. Evaluation, 14(1), 29-48.", {"size": 14}),
    ("  Koleros, A. & Mayne, J. (2017). Using Actor-Based Theories of Change. Journal of Development Effectiveness, 9(4), 439-460.", {"size": 14}),
    ("", {}),
    ("Contribution Analysis:", {"bold": True, "color": MID_BLUE}),
    ("  Mayne, J. (2012). Contribution Analysis: Coming of Age? Evaluation, 18(3), 270-280.", {"size": 14}),
    ("  Mayne, J. ILAC Brief on Contribution Analysis.", {"size": 14}),
    ("", {}),
    ("Behavioral Science:", {"bold": True, "color": MID_BLUE}),
    ("  Michie, S., van Stralen, M.M. & West, R. (2011). The Behaviour Change Wheel. Implementation Science, 6:42.", {"size": 14}),
    ("", {}),
    ("Early Childhood Development:", {"bold": True, "color": MID_BLUE}),
    ("  Britto, P. et al. (2017). Nurturing Care: Promoting Early Childhood Development. Lancet.", {"size": 14}),
    ("  Gertler, P. et al. (2014). Labor Market Returns to Jamaica Reach Up Program. Science, 344(6187).", {"size": 14}),
], num=49, note="Reference slide -- not presented in session.")

# ── Slide 50: Acknowledgments ──
make_title_slide(
    "EVALAC I -- Session 1\nIntroduction to M&E and Theory of Change",
    "CLEAR-LAB  |  IDB  |  OVE\n\nAcknowledgments: Course design informed by BID/OVE feedback cycle.\n"
    "Case study: Creciendo Juntos (fictional, based on evidence from LAC ECD programs).\n"
    "COM-B framework: Michie, van Stralen & West (2011).",
    num=50
)

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════

out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "S1_Final_v02.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
